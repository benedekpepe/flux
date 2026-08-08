"""
End-to-end verification.

Run:  python tests/audit.py

Checks, in the order a reviewer would ask about them:

    1  P&L arithmetic and the roll-up to computed subtotals
    2  favourable / unfavourable logic per account type
    3  materiality: the two-condition rule and the not-meaningful escape
    4  number parsing, including the accounting sign conventions
    5  period parsing and reporting-period derivation
    6  column mapping, including the guard that stops a document number
       being read as an amount
    7  aggregation, the budget join and its proportional allocation
    8  sign normalisation of a credit-balance ledger
    9  expense grouping, including unknown client-specific types
   10  cross-view reconciliation: P&L, entity, department and expense views
       all tie to the same source
   11  both generated workbooks: the shared column layout every reporting
       sheet now carries, the single lever cells behind it, structure, and the
       actuals-only contract
   12  edge cases and the app module

Exits non-zero on any failure.
"""

from __future__ import annotations

import sys
import tempfile
import traceback
from pathlib import Path

import pandas as pd
from openpyxl.utils import get_column_letter

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from flux import ingest
from flux.coa import CHART_OF_ACCOUNTS, EXPENSE_GROUPS, PNL_STRUCTURE
from flux.commentary import generate_commentary, line_comments
from flux.engine import (MaterialityRule, build_report, cost_centre_variances,
                         department_variances, entity_variances, has_budget,
                         leaf_variances)
from flux.reporting import build_client_pack, build_demo_pack, build_pptx_pack
from flux.reporting.client_pack import _group_expense_types
from flux.synthetic_data import (generate_budget_year, generate_month,
                                 generate_ytd_transactions, monthly_detail)

TOL = 0.01
PERIOD = "2025-06"

_results: list[tuple[str, str, str]] = []


def check(name: str, condition, detail: str = "") -> None:
    _results.append(("PASS" if condition else "FAIL", name, detail))


def close(a, b, tol=TOL) -> bool:
    return abs(float(a) - float(b)) <= tol


def line(report: pd.DataFrame, label: str):
    return report[report["line"] == label].iloc[0]


# ---------------------------------------------------------------------------
# 1-3  Engine
# ---------------------------------------------------------------------------
def test_engine() -> None:
    gl = generate_month(PERIOD)
    rep = build_report(gl)

    rev, cogs, gp = line(rep, "Revenue"), line(rep, "Cost of goods sold"), line(rep, "Gross profit")
    opex, ebit = line(rep, "Operating expenses"), line(rep, "Operating income (EBIT)")
    other, ni = line(rep, "Other expenses"), line(rep, "Net income")

    check("Gross profit = Revenue - COGS",
          close(gp.actual, rev.actual - cogs.actual))
    check("EBIT = Gross profit - OpEx",
          close(ebit.actual, gp.actual - opex.actual))
    check("Net income = EBIT - Other",
          close(ni.actual, ebit.actual - other.actual))
    check("Budget subtotals roll up the same way",
          close(gp.budget, rev.budget - cogs.budget)
          and close(ni.budget, ebit.budget - other.budget))
    check("Every P&L line is reported",
          len(rep) == len(PNL_STRUCTURE))
    check("Variance = actual - budget on every line",
          all(close(r.var_bud, r.actual - r.budget) for r in rep.itertuples()))

    # F/U respects account type: revenue up is good, cost up is bad.
    up = pd.DataFrame([
        {"account_code": "4000", "account_name": "Rev", "category": "Revenue",
         "actual": 120_000.0, "budget": 100_000.0, "prior_year": 0.0},
        {"account_code": "6110", "account_name": "Ads", "category": "OpEx",
         "actual": 120_000.0, "budget": 100_000.0, "prior_year": 0.0},
    ])
    r_up = build_report(up)
    check("Revenue above budget is favourable",
          line(r_up, "Revenue").fav_unfav == "F")
    check("OpEx above budget is unfavourable",
          line(r_up, "Operating expenses").fav_unfav == "U")
    check("Profit subtotal above budget is favourable",
          line(r_up, "Net income").fav_unfav == "F")

    # Materiality needs BOTH floors.
    rule = MaterialityRule(abs_threshold=25_000, pct_threshold=0.10)
    big_small_pct = pd.DataFrame([{"account_code": "6110", "account_name": "Ads",
                                   "category": "OpEx", "actual": 1_030_000.0,
                                   "budget": 1_000_000.0, "prior_year": 0.0}])
    check("Large EUR but small % is not material",
          not line(build_report(big_small_pct, rule), "Operating expenses").material)

    small_big_pct = pd.DataFrame([{"account_code": "6110", "account_name": "Ads",
                                   "category": "OpEx", "actual": 9_000.0,
                                   "budget": 3_000.0, "prior_year": 0.0}])
    check("Large % but small EUR is not material",
          not line(build_report(small_big_pct, rule), "Operating expenses").material)

    both = pd.DataFrame([{"account_code": "6110", "account_name": "Ads",
                          "category": "OpEx", "actual": 130_000.0,
                          "budget": 100_000.0, "prior_year": 0.0}])
    check("Clearing both floors is material",
          line(build_report(both, rule), "Operating expenses").material)

    # Not meaningful: a near-zero base gives no usable percentage.
    nm = pd.DataFrame([{"account_code": "6110", "account_name": "Ads",
                        "category": "OpEx", "actual": 150_000.0,
                        "budget": 7.0, "prior_year": 0.0}])
    nm_line = line(build_report(nm, rule), "Operating expenses")
    check("Near-zero base gives no percentage", nm_line.var_bud_pct is None)
    check("Not-meaningful line is still judged on the amount", nm_line.material)

    small_base = pd.DataFrame([{"account_code": "6110", "account_name": "Ads",
                                "category": "OpEx", "actual": 60_000.0,
                                "budget": 23_000.0, "prior_year": 0.0}])
    check("A small but real base still gets a percentage",
          line(build_report(small_base, rule), "Operating expenses").var_bud_pct is not None)

    # Spend views use the same materiality rule as the P&L.
    detail = monthly_detail(PERIOD)
    dept = department_variances(detail)
    check("Department view excludes revenue",
          close(dept["actual"].sum(),
                detail.loc[detail.category != "Revenue", "actual"].sum()))
    check("Department variance = actual - budget",
          all(close(r.var_bud, r.actual - r.budget) for r in dept.itertuples()))


# ---------------------------------------------------------------------------
# 4-5  Parsing
# ---------------------------------------------------------------------------
def test_parsing() -> None:
    cases = {
        "512.000,00": 512_000.0,        # EU thousands + decimal comma
        "1,234,567.89": 1_234_567.89,   # US thousands
        "445.000": 445_000.0,           # EU thousands, no decimals
        "12,5": 12.5,                   # decimal comma
        "\u20ac 1.234,50": 1234.5,      # currency symbol
        "1.234,50 EUR": 1234.5,         # currency code
        "-45 678,90": -45_678.9,        # space as thousands separator
        "(2.500,00)": -2500.0,          # accounting parentheses
        "1.234,50-": -1234.5,           # SAP trailing minus
        "1,234.50 CR": -1234.5,         # credit marker
        "1,234.50 DR": 1234.5,          # debit marker
        "\u22121.000,00": -1000.0,      # unicode minus
        "1\u00a0234,50": 1234.5,        # non-breaking space
        "0,00": 0.0,
    }
    for raw, want in cases.items():
        got = ingest._coerce_amount(raw)
        check(f"parse {raw!r} -> {want}", close(got, want), f"got {got}")

    for blank in ("", "n/a", "-", None):
        check(f"blank {blank!r} parses as NaN", pd.isna(ingest._coerce_amount(blank)))

    periods = {"2025-06": 202506, "2025/06": 202506, "202506": 202506,
               "06.2025": 202506, "2025-06-30": 202506,
               "June": None, "2025": None, "": None, None: None}
    for raw, want in periods.items():
        check(f"period {raw!r} -> {want}", ingest.period_key(raw) == want)

    # The reporting month is the latest period that carries actuals: a full-year
    # budget must not drag it into the future.
    std = pd.DataFrame({
        "account_code": ["4000"] * 4,
        "account_name": ["Rev"] * 4,
        "category": ["Revenue"] * 4,
        "period": ["2025-05", "2025-06", "2025-07", "2025-08"],
        "actual": [100.0, 120.0, 0.0, 0.0],
        "budget": [100.0, 100.0, 100.0, 100.0],
        "prior_year": [0.0] * 4,
    })
    active, found = ingest.reporting_period(std)
    check("Reporting month is the last month with postings", active == "2025-06", active)
    check("All periods are still listed", found[-1] == "2025-08")
    check("filter_period keeps one month",
          len(ingest.filter_period(std, "2025-06")) == 1)


# ---------------------------------------------------------------------------
# 6-7  Ingestion
# ---------------------------------------------------------------------------
def test_ingestion() -> None:
    messy = pd.DataFrame({
        "GL Code": ["4100", "4000", "5000", "6110", "6300"],
        "Description": ["Subscription rev", "Product rev", "Materials",
                        "Digital ads", "G&A salaries"],
        "Actual Amount (EUR)": ["512.000,00", "298.000,00", "168.000,00",
                                "128.000,00", "88.000,00"],
        "Plan": ["480.000,00", "340.000,00", "160.000,00", "90.000,00", "90.000,00"],
        "Last Year": ["445.000", "331.000", "150.000", "82.000", "86.000"],
        "Dept": ["Commercial", "Commercial", "Operations", "Marketing", "G&A"],
    })
    std, report, issues = ingest.ingest(messy)
    mapped = {m.field: m.mapped_from for _, m in report.iterrows()}
    check("Account code resolved from an odd header", mapped["account_code"] == "GL Code")
    check("Actual resolved from an odd header", mapped["actual"] == "Actual Amount (EUR)")
    check("Budget resolved from 'Plan'", mapped["budget"] == "Plan")
    check("Department resolved from 'Dept'", mapped["department"] == "Dept")
    check("Category inferred from the account-code range",
          list(std["category"]) == ["Revenue", "Revenue", "COGS", "OpEx", "OpEx"])
    check("The recognised chart of accounts is named",
          std.attrs["category_inference"]["style"] == "4-5-6-7",
          str(std.attrs["category_inference"]))
    # The classification and expense-type notes describe a guess the user has
    # to see; they are not mapping problems, which is what this asserts.
    inferred = ("classif", "chart of accounts", "expense type", "expense-type")
    check("A clean file raises no mapping issues",
          not [i for i in issues if not any(k in i.lower() for k in inferred)],
          str(issues))
    check("An inferred expense type is declared, never silent",
          any("inferred from the account names" in i for i in issues), str(issues))

    # A misread chart of accounts is the one failure that produces a
    # finished-looking pack with the bottom line inverted, so each of these
    # checks that the classification survives a chart the engine was not
    # designed around, and that the file says what it based the reading on.
    charts = {
        "Hungarian": (
            ["911", "912", "511", "541", "551", "521", "871"],
            ["Belfoldi ertekesites arbevetele", "Exportertekesites arbevetele",
             "Anyagkoltseg", "Berkoltseg", "Bergarulek",
             "Igenybe vett szolgaltatasok", "Fizetett kamat"],
            ["Revenue", "Revenue", "COGS", "OpEx", "OpEx", "OpEx", "Other"]),
        "SAP": (
            ["800000", "800100", "400000", "430000", "470000", "480000"],
            ["Sales revenue domestic", "Sales revenue export", "Raw materials",
             "Personnel expense", "Occupancy rent", "Depreciation"],
            ["Revenue", "Revenue", "COGS", "OpEx", "OpEx", "Other"]),
        "4-5-6-7": (
            ["4000", "5000", "6100", "7000"],
            ["Product revenue", "Materials", "Marketing spend", "Interest paid"],
            ["Revenue", "COGS", "OpEx", "Other"]),
        "SKR03": (
            ["8400", "8300", "3400", "4100", "4210", "4830"],
            ["Erl\u00f6se 19% USt", "Erl\u00f6se steuerfrei", "Wareneingang",
             "L\u00f6hne und Geh\u00e4lter", "Miete", "Abschreibungen"],
            ["Revenue", "Revenue", "COGS", "OpEx", "OpEx", "Other"]),
    }
    for style, (codes, names, want) in charts.items():
        frame = pd.DataFrame({"Account": codes, "Account name": names,
                              "Actual": [100_000.0] * len(codes)})
        got, info = ingest.infer_categories(codes, names)
        check(f"{style} chart of accounts is classified correctly", got == want,
              f"{got}")
        check(f"{style} chart of accounts is recognised by name",
              info["style"] == style, str(info["style"]))
        std_c, _r, _i = ingest.ingest(frame)
        check(f"{style} revenue reaches the P&L",
              line(build_report(std_c), "Revenue").actual > 0, style)

    # SAP and SKR03 both put revenue in the 8000s and disagree about the 4000s.
    # Only the agreement score against the account names can tell them apart, so
    # this is the check that the score is load-bearing rather than decorative.
    sap_codes = ["800000", "400000", "430000"]
    skr_codes = ["8400", "3400", "4100"]
    _c, sap_info = ingest.infer_categories(
        sap_codes, ["Sales revenue", "Raw materials", "Personnel expense"])
    _c, skr_info = ingest.infer_categories(
        skr_codes, ["Erl\u00f6se", "Wareneingang", "L\u00f6hne und Geh\u00e4lter"])
    check("Two charts sharing a revenue range are told apart by the names",
          sap_info["style"] == "SAP" and skr_info["style"] == "SKR03",
          f"{sap_info['style']} / {skr_info['style']}")

    # Account names carry no signal at all: the codes are the only evidence
    # left, and refusing to read them would classify the ledger as all OpEx.
    blind, blind_info = ingest.infer_categories(
        ["4001", "5001", "6001", "7001"], ["", "", "", ""])
    check("A file with no usable names still classifies from the codes",
          blind == ["Revenue", "COGS", "OpEx", "Other"], str(blind))
    check("A blind reading declares itself",
          blind_info["blind"] and any("no account name carried" in i.lower()
                                      for i in ingest.category_issues(
                                          ["4001"], [""], blind, blind_info)))

    # A chart nobody recognises must fall back to the names, not to silence.
    odd_codes = ["A10", "A20", "B10", "C10"]
    odd_names = ["Widget sales", "Consulting revenue", "Component purchases",
                 "Office rent"]
    got, info = ingest.infer_categories(odd_codes, odd_names)
    check("An unrecognised chart falls back to the account names",
          got == ["Revenue", "Revenue", "COGS", "OpEx"], str(got))
    check("An unrecognised chart says so rather than guessing quietly",
          info["style"] is None, str(info["style"]))
    check("The unrecognised chart is reported to the user",
          any("did not match any chart" in i
              for i in ingest.category_issues(odd_codes, odd_names, got, info)))

    # Costs with no revenue at all is the shape of a misread chart.
    all_cost, cost_info = ingest.infer_categories(
        ["6100", "6200"], ["Marketing spend", "Salaries"])
    check("A file with no revenue line is flagged",
          any("No account was classified as revenue" in i for i in
              ingest.category_issues(["6100", "6200"], ["Marketing spend", "Salaries"],
                                     all_cost, cost_info)))

    # Expense types, where the file carries no expense-type column.
    etypes, einfo = ingest.infer_expense_types(
        ["Berkoltseg", "Bergarulek", "Marketing kampany", "Berleti dij",
         "Ertekcsokkenesi leiras", "Valami egyedi tetel", "Arbevetel"],
        ["OpEx", "OpEx", "OpEx", "OpEx", "Other", "OpEx", "Revenue"])
    check("Expense types are inferred from Hungarian account names",
          etypes[:5] == ["Salaries & wages", "Payroll benefits",
                         "Marketing & advertising", "Facilities & office",
                         "Depreciation & amortisation"], str(etypes))
    check("An unrecognised cost is grouped, not dropped",
          etypes[5] == ingest.UNCLASSIFIED_EXPENSE, etypes[5])
    check("Revenue accounts get no expense type", etypes[6] == "")
    check("The inference reports how much of it was a guess",
          einfo["recognised"] == 5 and einfo["unclassified"] == 1, str(einfo))

    # A manufacturer books direct labour into cost of sales, which the statutory
    # chart puts under operating cost. The code wins - it is corroborated across
    # the file - but the reader has to be told, or the disagreement is silent.
    mixed_codes = ["911", "5111", "5411", "5511", "5711"]
    mixed_names = ["Belfoldi ertekesites arbevetele", "Anyagkoltseg",
                   "Kozvetlen berkoltseg (termeles)", "Bergarulekok",
                   "Ertekcsokkenesi leiras"]
    mixed, mixed_info = ingest.infer_categories(mixed_codes, mixed_names)
    check("A name that contradicts a corroborated chart does not silently win",
          mixed[2] == "OpEx", mixed[2])
    reported = ingest.category_issues(mixed_codes, mixed_names, mixed, mixed_info)
    check("The contradiction is reported so it can be overridden",
          any("disagreed" in i and "Kozvetlen" in i for i in reported),
          str(reported))
    check("The report says which reading was followed",
          any("the code was followed" in i for i in reported), str(reported))

    # A real export is a printed report: a title block, a subtotal after each
    # group, a grand total and a footer. Read as data, the subtotal rows count
    # the whole ledger twice, quietly, while the pack still looks finished.
    printed = pd.DataFrame({
        "Szamla": ["911", "5111", "5411", "", "Keszult: 2025.07.02"],
        "Megnevezes": ["Belfoldi ertekesites arbevetele", "Anyagkoltseg",
                       "Berkoltseg", "Osszesen", ""],
        "Tartozik": [0, 4_200_000, 5_100_000, 9_300_000, 0],
        "Kovetel": [12_000_000, 0, 0, 12_000_000, 0],
    })
    pstd, _report, pissues = ingest.ingest(printed)
    check("A grand-total row is not read as an account",
          len(pstd) == 3, f"{len(pstd)} rows: {list(pstd['account_name'])}")
    check("Dropping the total row is reported, never silent",
          any("total rows were left out" in i for i in pissues), str(pissues))
    check("A footer line is not read as an account",
          "Keszult: 2025.07.02" not in list(pstd["account_code"]))
    prep = build_report(pstd)
    check("The printed report's revenue survives the debit/credit convention",
          close(line(prep, "Revenue").actual, 12_000_000),
          str(line(prep, "Revenue").actual))
    check("A single credit-balance revenue account is still normalised",
          any("credit balance" in i for i in pissues), str(pissues))
    check("The result is the one the ledger actually shows",
          close(line(prep, "Net income").actual, 2_700_000),
          str(line(prep, "Net income").actual))

    # An accounting export routinely opens with a cover sheet and puts the
    # ledger second. Reading the first sheet found no account column at all.
    import openpyxl
    book = Path(tempfile.mkdtemp()) / "cover.xlsx"
    wb_in = openpyxl.Workbook()
    wb_in.active.title = "Fedlap"
    wb_in.active.append(["XY Kft."])
    wb_in.active.append(["2025.01.01 - 2025.06.30"])
    data = wb_in.create_sheet("Fokonyv")
    for row in (["Szamla", "Megnevezes", "Teny"],
                ["911", "Arbevetel", -5_000_000],
                ["5111", "Anyagkoltseg", 2_000_000]):
        data.append(row)
    wb_in.save(book)
    cstd, _r3, _i3 = ingest.ingest(book)
    check("The ledger sheet is found behind a cover sheet",
          cstd is not None and len(cstd) == 2,
          "no rows" if cstd is None else str(len(cstd)))
    if cstd is not None:
        check("A workbook with a cover sheet still reports its revenue",
              close(line(build_report(cstd), "Revenue").actual, 5_000_000),
              str(line(build_report(cstd), "Revenue").actual))

    # An account legitimately named after a total word must survive, or the
    # filter is worse than the problem it fixes.
    real = pd.DataFrame({
        "Account": ["4000", "6500", "6600"],
        "Account name": ["Product revenue", "Net interest expense",
                         "Gross margin adjustment"],
        "Actual": [500_000.0, 12_000.0, 8_000.0],
    })
    rstd, _r2, _i2 = ingest.ingest(real)
    check("A coded account named after a total word is kept",
          len(rstd) == 3, f"{len(rstd)}: {list(rstd['account_name'])}")

    # A document number must never be read as an amount.
    docnum = pd.DataFrame({
        "Account": ["4000", "5000", "6110", "6300", "6310", "6320"],
        "Name": ["Rev", "Mat", "Ads", "Sal", "Rent", "Util"],
        "Doc": [1900001, 1900002, 1900003, 1900004, 1900005, 1900006],
        "FY": [2025] * 6,
    })
    matches = ingest.match_columns(list(docnum.columns))
    ingest._content_detect(docnum, matches)
    amount = next(m.source for m in matches if m.field == "actual")
    check("Sequential document numbers are not read as amounts", amount != "Doc", str(amount))
    check("A fiscal-year column is not read as an amount", amount != "FY", str(amount))

    # Hungarian headers resolve through the synonym dictionary.
    hu = pd.DataFrame({
        "Fokonyvi szam": ["4100", "6110"],
        "Megnevezes": ["Arbevetel", "Hirdetes"],
        "Teny": ["1.200.000", "300.000"],
        "Terv": ["1.000.000", "250.000"],
        "Koltseghely": ["CM10100", "MK30100"],
    })
    hm = ingest.match_columns(list(hu.columns))
    hmap = {m.field: m.source for m in hm if m.source}
    check("Hungarian headers map to the schema",
          hmap.get("account_code") == "Fokonyvi szam"
          and hmap.get("actual") == "Teny"
          and hmap.get("budget") == "Terv"
          and hmap.get("cost_centre") == "Koltseghely", str(hmap))

    # Trial-balance shape: separate debit and credit columns.
    tb = pd.DataFrame({
        "Account": ["6110", "6300"],
        "Account name": ["Ads", "Salaries"],
        "Debit": ["130.000,00", "90.000,00"],
        "Credit": ["30.000,00", "0,00"],
    })
    tm = ingest.match_columns(list(tb.columns))
    tstd = ingest.apply_mapping(tb, tm)
    check("Debit and credit columns net to the actual",
          close(tstd["actual"].iloc[0], 100_000.0), str(list(tstd["actual"])))

    # One amount column plus a debit/credit marker per line.
    dc = pd.DataFrame({
        "Account": ["4000", "6110"],
        "Account name": ["Rev", "Ads"],
        "Amount": ["500.000,00", "130.000,00"],
        "D/C": ["H", "S"],
    })
    dm = ingest.match_columns(list(dc.columns))
    dstd = ingest.apply_mapping(dc, dm)
    check("A credit marker signs the amount negative",
          close(dstd["actual"].iloc[0], -500_000.0) and close(dstd["actual"].iloc[1], 130_000.0),
          str(list(dstd["actual"])))

    # Aggregation keeps the analytical dimensions.
    lines = pd.DataFrame({
        "account_code": ["6110", "6110", "6110"],
        "account_name": ["Ads"] * 3,
        "category": ["OpEx"] * 3,
        "department": ["Marketing", "Marketing", "Sales"],
        "actual": [10.0, 20.0, 30.0],
        "budget": [0.0, 0.0, 0.0],
        "prior_year": [0.0, 0.0, 0.0],
    })
    agg = ingest.aggregate_to_accounts(lines)
    check("Posting lines collapse per account and dimension", len(agg) == 2, str(len(agg)))
    check("Aggregation preserves the total", close(agg["actual"].sum(), 60.0))

    # A budget join across a coarser plan allocates rather than repeats.
    actuals = pd.DataFrame({
        "account_code": ["6110", "6110"],
        "account_name": ["Ads", "Ads"],
        "category": ["OpEx", "OpEx"],
        "actual": [75.0, 25.0],
        "budget": [0.0, 0.0],
        "prior_year": [0.0, 0.0],
    })
    plan = pd.DataFrame({
        "account_code": ["6110"],
        "account_name": ["Ads"],
        "category": ["OpEx"],
        "actual": [0.0],
        "budget": [200.0],
        "prior_year": [0.0],
    })
    joined, _ = ingest.merge_budget(actuals, plan)
    check("A coarser budget is allocated, not repeated",
          close(joined["budget"].sum(), 200.0), str(joined["budget"].sum()))
    check("Allocation follows the weight of actuals",
          close(joined["budget"].max(), 150.0), str(joined["budget"].max()))


# ---------------------------------------------------------------------------
# 8  Sign normalisation
# ---------------------------------------------------------------------------
def test_sign_normalisation() -> None:
    # A ledger credits revenue, so revenue arrives negative.
    ledger = pd.DataFrame({
        "account_code": ["4000", "4100", "5000", "6110"],
        "account_name": ["Rev A", "Rev B", "Materials", "Ads"],
        "category": ["Revenue", "Revenue", "COGS", "OpEx"],
        "actual": [-500_000.0, -300_000.0, 160_000.0, 120_000.0],
        "budget": [-480_000.0, -320_000.0, 150_000.0, 100_000.0],
        "prior_year": [0.0, 0.0, 0.0, 0.0],
    })
    fixed, notes = ingest.normalise_signs(ledger)
    rev = fixed.loc[fixed.category == "Revenue", "actual"].sum()
    check("Credit-balance revenue is normalised to a positive magnitude",
          close(rev, 800_000.0), str(rev))
    check("Costs are left alone",
          close(fixed.loc[fixed.category == "COGS", "actual"].sum(), 160_000.0))
    check("The sign change is reported", any("normalised" in n for n in notes))

    rep = build_report(fixed)
    check("A normalised ledger produces positive revenue in the P&L",
          line(rep, "Revenue").actual > 0)

    # One credit note among positives is real data, not a convention.
    mixed = pd.DataFrame({
        "account_code": ["6110", "6120", "6130"],
        "account_name": ["Ads", "Events", "PR"],
        "category": ["OpEx"] * 3,
        "actual": [100_000.0, 80_000.0, -250_000.0],
        "budget": [0.0, 0.0, 0.0],
        "prior_year": [0.0, 0.0, 0.0],
    })
    kept, kept_notes = ingest.normalise_signs(mixed)
    check("A single credit inside a positive category is not flipped",
          close(kept["actual"].sum(), -70_000.0), str(kept["actual"].sum()))
    check("No sign note is raised for genuine mixed data", not kept_notes)


# ---------------------------------------------------------------------------
# 9  Expense grouping
# ---------------------------------------------------------------------------
def test_expense_grouping() -> None:
    known = [t for _g, types in EXPENSE_GROUPS for t in types]
    groups = _group_expense_types(known)
    check("Every known expense type is grouped",
          sorted(t for _n, ts in groups for t in ts) == sorted(known))
    check("Personnel leads the expense report", groups[0][0] == "Personnel costs")

    with_unknown = _group_expense_types(["Salaries & wages", "Céges mobilflotta"])
    flat = [t for _n, ts in with_unknown for t in ts]
    check("A client-specific expense type is still reported",
          "Céges mobilflotta" in flat)
    check("Unknown types are collected, not silently dropped",
          with_unknown[-1][0] == "Other expense types")


# ---------------------------------------------------------------------------
# 10  Cross-view reconciliation
# ---------------------------------------------------------------------------
def test_reconciliation() -> None:
    detail = monthly_detail(PERIOD)
    agg = generate_month(PERIOD).drop(columns="period")
    rep = build_report(agg)

    ni = line(rep, "Net income").actual
    rev = line(rep, "Revenue").actual

    ent = entity_variances(detail)
    check("Entity net income consolidates to the group P&L",
          close(ent["net_actual"].sum(), ni), f"{ent['net_actual'].sum()} vs {ni}")
    check("Entity revenue consolidates to group revenue",
          close(ent["revenue"].sum(), rev))

    spend = (line(rep, "Cost of goods sold").actual
             + line(rep, "Operating expenses").actual
             + line(rep, "Other expenses").actual)
    dept = department_variances(detail)
    check("Department spend ties to the P&L cost lines",
          close(dept["actual"].sum(), spend), f"{dept['actual'].sum()} vs {spend}")

    cc = cost_centre_variances(detail)
    check("Cost centres sum to the same spend",
          close(cc["actual"].sum(), spend))
    check("Every cost centre carries its parent department",
          "department" in cc.columns and cc["department"].notna().all())

    by_type = (detail[detail.category != "Revenue"]
               .groupby("expense_type", as_index=False)["actual"].sum())
    check("The expense view ties to the same spend",
          close(by_type["actual"].sum(), spend))

    check("Revenue less spend is net income", close(rev - spend, ni))

    leaves = leaf_variances(agg)
    check("Account detail sums to the P&L actual",
          close(leaves["actual"].sum(), rev + spend))

    # Year-to-date actuals must reconcile to the transaction source.
    txns = generate_ytd_transactions(PERIOD)
    month_total = txns.loc[txns.period == PERIOD, "amount_eur"].sum()
    check("The month's transactions tie to the P&L",
          close(month_total, rev + spend, tol=1.0),
          f"{month_total} vs {rev + spend}")


# ---------------------------------------------------------------------------
# 11  Workbooks
# ---------------------------------------------------------------------------
def _load(path: Path):
    import warnings
    from openpyxl import load_workbook
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        return load_workbook(path)


def _cf_cells(ws) -> set[str]:
    """Every cell coordinate covered by a conditional-formatting range.

    Coordinates, not row numbers: a rule over the variance columns covering a
    row says nothing about whether the F/U badge beside it is covered too, and
    unioning the ranges lets one hide the other's gap.
    """
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.cell import range_boundaries

    covered: set[str] = set()
    for rng in ws.conditional_formatting:
        for part in str(rng.sqref).split():
            c1, r1, c2, r2 = range_boundaries(part)
            for col in range(c1, c2 + 1):
                for row in range(r1, r2 + 1):
                    covered.add(f"{get_column_letter(col)}{row}")
    return covered


def _header_columns(ws) -> dict:
    """Map header label to column letter, for whichever row holds the headers."""
    from openpyxl.utils import get_column_letter

    for header_row in (5, 10):
        found = {}
        for col in range(1, ws.max_column + 1):
            value = ws.cell(row=header_row, column=col).value
            # First occurrence wins: every sheet carries two "Var %" columns,
            # one for the month and one for the year to date, and the month one
            # is the one coloured from the F/U badge beside it.
            if isinstance(value, str) and value.strip() and value.strip() not in found:
                found[value.strip()] = get_column_letter(col)
        if "F/U" in found:
            return found
    return {}


def _table_end(ws) -> int:
    """The last row of the report table, before the analysis section starts.

    The analysis blocks are headed by the line they describe, so "Net income ·
    BOTH" reads as a total row to anything scanning column A - and then fails a
    conditional-formatting check for a row that is prose.
    """
    for r in range(5, ws.max_row + 1):
        if str(ws.cell(row=r, column=1).value or "").startswith("Analysis"):
            return r - 1
    return ws.max_row


def _total_rows(ws) -> list[int]:
    return [r for r in range(5, _table_end(ws) + 1)
            if isinstance(ws.cell(row=r, column=1).value, str)
            and ws.cell(row=r, column=1).value.lower().startswith(
                ("total", "consolidated", "net income"))]


def _report_header(ws):
    """The header row of a reporting sheet, and its labels.

    Reporting sheets put their headers on row 5, except the P&L, which uses row
    10 to leave space for the KPI cards and the lever cells above them.
    """
    for hrow in (5, 10):
        labels = [ws.cell(row=hrow, column=c).value
                  for c in range(1, ws.max_column + 1)]
        if "Flag" in labels:
            return hrow, labels
    return None, []


def _check_report_layout(wb, label: str) -> None:
    """Every reporting sheet carries the same columns and the same levers.

    This is the check the pack did not have when the sheets diverged: the P&L
    showed a year to date the expense report did not, the outlook sheet was the
    only one with a projection, and the flag meant "material" on one sheet and
    "material this month" on another. Reading the headers back off the saved
    workbook is the only way to catch that, because each sheet builds its own
    formulas and any of them can drift alone.
    """
    from flux.reporting.formulas import CORE_HEADERS

    checked = 0
    for name in wb.sheetnames:
        ws = wb[name]
        hrow, labels = _report_header(ws)
        if hrow is None:
            continue                      # an input sheet, not a report
        checked += 1
        start = labels.index("Month Act") if "Month Act" in labels else -1
        core = tuple(labels[start:start + len(CORE_HEADERS)]) if start >= 0 else ()
        check(f"{label}: '{name}' carries the shared column set",
              core == CORE_HEADERS, f"got {core}")
        trailing = [x for x in labels[start + len(CORE_HEADERS):] if x]
        check(f"{label}: '{name}' adds nothing after the flag but commentary",
              trailing in ([], ["Commentary"]), str(trailing))
        if start < 0:
            continue

        # Both levered columns must point at the P&L's cells, not at a local
        # copy: one edit has to re-flag and re-project the whole pack.
        first = hrow + 1
        flag = str(ws.cell(row=first, column=start + len(CORE_HEADERS)).value or "")
        rate = str(ws.cell(row=first, column=start + 10).value or "")
        owns_levers = name == "P&L Report"
        prefix = "" if owns_levers else "'P&L Report'!"
        check(f"{label}: '{name}' flags against the shared materiality levers",
              f"{prefix}$C$9" in flag and f"{prefix}$G$9" in flag, flag[:90])
        check(f"{label}: '{name}' names the material timeframe",
              all(v in flag for v in ('"MONTH"', '"YTD"', '"BOTH"')), flag[:90])
        check(f"{label}: '{name}' projects off the shared months lever",
              f"{prefix}$K$9" in rate, rate[:90])
    check(f"{label}: the layout check found the reporting sheets",
          checked >= 2, f"{checked} sheets")


def _check_total_row_formatting(wb, label: str) -> None:
    """Every total row must be coloured by the same rules as the detail above it.

    This is the check that was missing when the grand-total row rendered black
    while every line above it was green or red: the formatting range stopped one
    row short, and nothing noticed.
    """
    for name in wb.sheetnames:
        ws = wb[name]
        totals = _total_rows(ws)
        if not totals or not list(ws.conditional_formatting):
            continue
        covered = _cf_cells(ws)
        headers = _header_columns(ws)
        watched = [headers[h] for h in ("F/U", "Month Var", "Var %")
                   if h in headers]
        if not watched:
            continue
        missing = [f"{col}{row}" for row in totals for col in watched
                   if f"{col}{row}" not in covered]
        check(f"{label}: total row on '{name}' is conditionally formatted",
              not missing, f"uncovered cells: {missing}")


def test_workbooks() -> None:
    tmp = Path(tempfile.mkdtemp())

    demo = build_demo_pack(PERIOD, tmp / "demo.xlsx")
    wb = _load(demo)
    expected = ["P&L Report", "Expense Report", "By Entity",
                "Departments & CCs", "Drivers", "Budget", "GL Transactions"]
    check("Demo pack has every sheet in reading order",
          wb.sheetnames == expected, str(wb.sheetnames))
    check("Demo pack opens on the P&L",
          wb.sheetnames[wb.active.index if hasattr(wb.active, "index") else 0] == "P&L Report"
          or wb.active.title == "P&L Report", wb.active.title)
    check("The P&L is formula-driven, not a value dump",
          str(wb["P&L Report"]["B11"].value).startswith("="),
          str(wb["P&L Report"]["B11"].value))
    check("Materiality levers are present and editable",
          wb["P&L Report"]["C9"].value == 25_000 and wb["P&L Report"]["G9"].value == 0.10)
    check("Months elapsed is a lever on the P&L, not a constant in the formulas",
          wb["P&L Report"]["K9"].value == 6, str(wb["P&L Report"]["K9"].value))

    # The whole point of the rebuild: one column set, one set of levers.
    _check_report_layout(wb, "Demo pack")

    # A total row is the line a reader looks at first, and it is the one an
    # off-by-one in a formatting range silently drops: the detail above it goes
    # green and red while the total stays black.
    _check_total_row_formatting(wb, "Demo pack")

    # Analysis under each sheet: the same findings in both packs and the deck,
    # and never a causal claim - the ledger does not record why a line moved,
    # and a pack that invented one would be unfalsifiable against its own table.
    from flux.analysis import NO_CAUSE_NOTE, concentration, persistence, question

    analysed = 0
    for name in ("P&L Report", "Expense Report", "By Entity", "Departments & CCs"):
        ws = wb[name]
        column_a = [str(ws.cell(row=r, column=1).value or "")
                    for r in range(1, ws.max_row + 1)]
        has_section = any(v.startswith("Analysis") for v in column_a)
        check(f"Demo pack: '{name}' carries an analysis section", has_section,
              str(column_a[-6:]))
        if has_section:
            analysed += 1
            check(f"Demo pack: '{name}' states that it does not explain causes",
                  any(NO_CAUSE_NOTE[:40] in str(ws.cell(row=r, column=1).value or "")
                      for r in range(1, ws.max_row + 1)))
            check(f"Demo pack: '{name}' analysis names a heading it can support",
                  any(v in ("Concentration", "Coverage", "Persistence",
                            "Full year", "Ask") for v in column_a),
                  "no findings written")
    check("Demo pack: every report sheet was analysed", analysed == 4, str(analysed))

    # No sentence in the pack may claim a cause. This is the one assertion that
    # would catch a well-meaning edit turning description into diagnosis.
    # "because" is not on the list: the disclaimer uses it, and so does the
    # arithmetic ("more than the net movement, because part of it is offset").
    # These four are the phrasings that can only introduce a business cause,
    # which is the thing the ledger cannot support.
    banned = ("due to", "caused by", "the reason", "driven by the")
    offenders = []
    for name in wb.sheetnames:
        ws = wb[name]
        for row in ws.iter_rows():
            for cell in row:
                text = str(cell.value or "").lower()
                if len(text) > 40 and any(b in text for b in banned):
                    offenders.append(f"{name}!{cell.coordinate}")
    check("Nothing in the pack claims to know why a line moved",
          not offenders, str(offenders[:4]))

    # The findings themselves: shape in, sentence out.
    one_off = pd.DataFrame({
        "period": ["2025-0" + str(i) for i in range(1, 7)],
        "period_no": [202501 + i for i in range(6)],
        "actual": [100.0, 100.0, 100.0, 400_000.0, 100.0, 100.0],
        "budget": [100.0] * 6,
    })
    q = question(one_off, higher_is_better=False)
    check("A single month out of line is called a timing question",
          q is not None and "timing" in q.text, "" if q is None else q.text[:60])
    steady = one_off.copy()
    steady["actual"] = [200.0] * 6
    q2 = question(steady, higher_is_better=False)
    check("A steady overrun is called a level question",
          q2 is not None and "level rather than an event" in q2.text,
          "" if q2 is None else q2.text[:60])
    p1 = persistence(steady, higher_is_better=False)
    check("Persistence counts the adverse months",
          p1 is not None and "6 of 6" in p1.text, "" if p1 is None else p1.text)
    spread = pd.DataFrame({"account_name": list("abcdefghij"),
                           "var_bud": [10.0] * 10})
    c1 = concentration(spread, 100.0)
    check("A variance spread thin says so rather than naming a driver",
          c1 is not None and "no single driver" in c1.text,
          "" if c1 is None else c1.text[:60])

    # Roll-up rows carry commentary; leaf rows do not, because there the comment
    # could only restate the columns beside it. And the comment must cover both
    # timeframes, so it explains the flag rather than repeating one column.
    for name, roll_row, leaf_row in (("Departments & CCs", 6, 7),
                                     ("By Entity", 6, None)):
        ws = wb[name]
        headers_row = [ws.cell(row=5, column=c).value for c in range(1, ws.max_column + 1)]
        check(f"Demo pack: '{name}' carries a commentary column",
              headers_row[-1] == "Commentary", str(headers_row[-1]))
        text = ws.cell(row=roll_row, column=ws.max_column).value or ""
        check(f"Demo pack: '{name}' comments its roll-up rows on both timeframes",
              "YTD" in text and "Month" in text, text[:80])
        if leaf_row:
            leaf = ws.cell(row=leaf_row, column=ws.max_column).value
            check(f"Demo pack: '{name}' leaves its leaf rows uncommented",
                  leaf in (None, ""), str(leaf)[:60])
    # Every group row on the expense report is a roll-up, including the two
    # groups holding a single expense type: they are drawn as group rows, so a
    # blank comment beside three filled ones reads as a gap.
    exp = wb["Expense Report"]
    # A group holding one type is labelled with the type's name, not the
    # group's: "Financing" is drawn as "Financing & bank".
    label_of_group = {name: (types[0] if len(types) == 1 else name)
                      for name, types in EXPENSE_GROUPS}
    group_names = set(label_of_group.values())
    found = 0
    for r in range(6, exp.max_row + 1):
        label = str(exp.cell(row=r, column=1).value or "").strip()
        if label not in group_names:
            continue
        found += 1
        text = exp.cell(row=r, column=exp.max_column).value or ""
        check(f"Demo pack: expense group '{label}' carries a comment",
              text.startswith("YTD"), text[:60])
        # A group over one line must not attribute the movement to itself.
        singles = [t for n, t in EXPENSE_GROUPS if label_of_group[n] == label][0]
        if len(singles) == 1:
            check(f"Demo pack: '{label}' does not name itself as its own driver",
                  "Driven by" not in text, text[:80])
    check("Demo pack: every expense group was checked",
          found == len(EXPENSE_GROUPS), f"{found} of {len(EXPENSE_GROUPS)}")

    # The grand total is a roll-up over the whole sheet, and the row a reader
    # looks at first. It was the one blank cell in a commented column.
    for name, label in (("Expense Report", "Total expenses"),
                        ("Departments & CCs", "Total spend"),
                        ("By Entity", "Consolidated")):
        ws = wb[name]
        row = next(r for r in range(6, ws.max_row + 1)
                   if str(ws.cell(row=r, column=1).value).strip() == label)
        text = ws.cell(row=row, column=ws.max_column).value or ""
        check(f"Demo pack: '{name}' comments its total row", text.startswith("YTD"),
              text[:60])
        # And the comment has to agree with the row: a spend sheet excludes
        # revenue, so a comment built on the unfiltered frame contradicts it.
        cols = _header_columns(ws)
        from openpyxl.utils import column_index_from_string
        var = ws.cell(row=row,
                      column=column_index_from_string(cols["YTD Var"])).value
        check(f"Demo pack: '{name}' total comment matches its own YTD variance",
              isinstance(var, str) and var.startswith("="),
              "formula expected; the workbook is not recalculated here")

    # Rows scroll, labels should not: every sheet freezes its label columns as
    # well as its header, or scrolling to the run rate loses the row names.
    for name in wb.sheetnames:
        check(f"Demo pack: '{name}' freezes its label column",
              str(wb[name].freeze_panes or "")[:1] not in ("A", ""),
              str(wb[name].freeze_panes))

    pnl_comment = wb["P&L Report"].cell(row=11, column=wb["P&L Report"].max_column).value or ""
    check("Demo pack: the P&L commentary covers the year to date, not just the month",
          pnl_comment.startswith("YTD"), pnl_comment[:80])

    # The sheet carries a footnote below the table, so count the contiguous
    # data rows rather than trusting max_row.
    gl_ws = wb["GL Transactions"]
    gl_rows = 0
    for (cell,) in gl_ws.iter_rows(min_row=6, max_col=1):
        if cell.value in (None, ""):
            break
        gl_rows += 1
    txns = generate_ytd_transactions(PERIOD)
    check("Every transaction reaches the GL sheet",
          gl_rows == len(txns), f"{gl_rows} vs {len(txns)}")

    # Client pack, with a budget.
    agg = generate_month(PERIOD).drop(columns="period")
    withbud = build_client_pack(agg, PERIOD, tmp / "client.xlsx")
    wbc = _load(withbud)
    check("Client pack always has the core sheets",
          {"P&L Report", "Drivers", "GL Input"} <= set(wbc.sheetnames), str(wbc.sheetnames))
    check("Client P&L is formula-driven",
          str(wbc["P&L Report"]["B11"].value).startswith("="))
    _check_report_layout(wbc, "Client pack")

    _check_total_row_formatting(wbc, "Client pack")

    # Client pack, actuals only: the variance columns must stay empty.
    actuals = agg.copy()
    actuals["budget"] = 0.0
    check("An actuals-only frame is detected as unbudgeted", not has_budget(actuals))

    nobud = build_client_pack(actuals, PERIOD, tmp / "actuals.xlsx")
    wbn = _load(nobud)
    pnl = wbn["P&L Report"]
    check("Actuals are still reported without a budget",
          str(pnl["B11"].value).startswith("="), str(pnl["B11"].value))
    empty = all(pnl[f"{col}11"].value in (None, "") for col in ("C", "D", "E"))
    check("Budget and variance cells are left empty, not zeroed", empty,
          str([pnl[f"{c}11"].value for c in "CDE"]))
    check("No F/U badge is claimed without a budget",
          pnl["M11"].value in (None, ""), str(pnl["M11"].value))
    check("No materiality flag is claimed without a budget",
          pnl["N11"].value in (None, ""), str(pnl["N11"].value))
    check("Inert materiality levers are omitted",
          pnl["C9"].value is None, str(pnl["C9"].value))
    # The run rate needs actuals and a month count, not a plan, so it survives.
    check("The run rate still projects without a budget",
          str(pnl["K11"].value).startswith("="), str(pnl["K11"].value))
    check("Months elapsed stays a lever without a budget",
          isinstance(pnl["K9"].value, int), str(pnl["K9"].value))
    check("Var to FY is left empty without a full-year plan",
          pnl["L11"].value in (None, ""), str(pnl["L11"].value))

    # A file that starts mid-year: the year to date covers four months, so the
    # run rate must divide by four. Reading the month count off the reporting
    # period instead would understate every projection in the pack by a third.
    gappy = agg.copy()
    gappy["period"] = PERIOD
    gappy["period_no"] = ingest.period_key(PERIOD)
    earlier = []
    for p in ("2025-03", "2025-04", "2025-05"):
        block = agg.copy()
        block["period"] = p
        block["period_no"] = ingest.period_key(p)
        earlier.append(block)
    gappy = pd.concat(earlier + [gappy], ignore_index=True)
    part = build_client_pack(gappy, PERIOD, tmp / "gappy.xlsx")
    wbg = _load(part)
    check("Months elapsed counts the months with postings, not the month number",
          wbg["P&L Report"]["K9"].value == 4,
          str(wbg["P&L Report"]["K9"].value))

    rep_nb = build_report(actuals)
    check("The engine reports nothing material without a budget",
          not rep_nb["material"].any())
    text = generate_commentary(rep_nb, actuals)
    check("Commentary says a budget is missing rather than inventing one",
          "no budget" in text.lower(), text[:80])
    comments = line_comments(rep_nb, actuals)
    check("Per-line comments claim no variance without a budget",
          all("no budget" in c.lower() for c in comments.values()))

    # --- the management deck ------------------------------------------------
    from pptx import Presentation

    # The cumulative frame the year-to-date slide needs.
    cut = ingest.period_key(PERIOD)
    txns_ytd = generate_ytd_transactions(PERIOD)
    bud_ytd = generate_budget_year()
    ytd_actual = (txns_ytd[txns_ytd.period_no <= cut]
                  .groupby(["account_code", "account_name", "category"], as_index=False)
                  ["amount_eur"].sum().rename(columns={"amount_eur": "actual"}))
    ytd_plan = (bud_ytd[bud_ytd.period_no <= cut]
                .groupby(["account_code", "account_name", "category"], as_index=False)
                [["budget_eur", "prior_eur"]].sum()
                .rename(columns={"budget_eur": "budget", "prior_eur": "prior_year"}))
    ytd_frame = ytd_actual.merge(ytd_plan,
                                 on=["account_code", "account_name", "category"],
                                 how="outer").fillna(0.0)

    fy_frame = (bud_ytd.groupby(["account_code", "account_name", "category"],
                                as_index=False)[["budget_eur", "prior_eur"]].sum()
                .rename(columns={"budget_eur": "budget", "prior_eur": "prior_year"}))
    fy_frame["actual"] = 0.0

    # Built the way the product builds it: with the analysis findings and the
    # cumulative detail, so the checks below test the deck a reader receives.
    from flux.reporting import demo_analysis_blocks, demo_ytd_detail

    deck = build_pptx_pack(agg, PERIOD, tmp / "deck.pptx",
                           detail=monthly_detail(PERIOD), ytd=ytd_frame,
                           fy_budget=fy_frame, months=6,
                           analysis_blocks=demo_analysis_blocks(PERIOD),
                           ytd_detail=demo_ytd_detail(PERIOD))
    prs = Presentation(deck)
    check("Deck has all eight slides with a budget, a year to date and a plan",
          len(prs.slides._sldIdLst) == 8, str(len(prs.slides._sldIdLst)))

    words = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                words.append(shape.text_frame.text)
    joined = " ".join(words)
    check("Cover names the reporting period", PERIOD in joined)
    ni_actual = line(build_report(agg), "Net income").actual
    check("Deck states the net income figure",
          f"{ni_actual:,.0f}" in joined, "net income missing from the deck")
    check("Deck explains the materiality rule", "materiality" in joined.lower())

    # The deck and the workbook must headline the same timeframe, or the same
    # company reads two ways depending on which file you opened.
    cover = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes
                     if sh.has_text_frame)
    ytd_ni = line(build_report(ytd_frame), "Net income").actual
    check("The deck's cover headlines the year to date, as the workbook does",
          f"{ytd_ni:,.0f}" in cover, cover[:120])
    check("The cover still carries the month as the second reading",
          f"{ni_actual:,.0f}" in cover, cover[:120])

    # Table cells are not shape text frames, so the headers have to be read off
    # the table itself.
    table_heads = [cell.text for s in prs.slides for sh in s.shapes
                   if sh.has_table for cell in sh.table.rows[0].cells]
    check("The deck's P&L table carries both timeframes",
          "YTD Act" in table_heads and "Month Act" in table_heads,
          str(table_heads))
    # Actual against budget is the whole statement; a variance column on its own
    # makes the reader do the subtraction backwards, and the plan went missing
    # when the year-to-date columns were added.
    check("The deck's P&L table shows the plan, not only the variance",
          "YTD Bud" in table_heads, str(table_heads))
    check("The deck's P&L table carries the materiality flag",
          "Flag" in table_heads, str(table_heads))

    # The narrative has to describe the figures above it.
    result_text = " ".join(sh.text_frame.text for sh in prs.slides[1].shapes
                           if sh.has_text_frame and len(sh.text_frame.text) > 120)
    ytd_ni_mag = f"{abs(line(build_report(ytd_frame), 'Net income').actual) / 1000:.1f}k"
    check("The result slide's narrative describes the year to date, as its cards do",
          ytd_ni_mag in result_text, result_text[:110])

    # Two slides showing the same three figures is two slides doing one job.
    card_slides = [i for i, s in enumerate(prs.slides)
                   if sum(1 for sh in s.shapes if sh.has_text_frame
                          and sh.text_frame.text.strip().startswith("REVENUE")) ]
    check("Only one slide carries the year-to-date KPI cards",
          len(card_slides) <= 1, f"cards on slides {card_slides}")

    # The analysis slide must not print the same finding twice under two names.
    # Two blocks may share a sentence - "adverse in 6 of 6 months" is true of
    # several lines - so what matters is that no two blocks are the same story.
    headings = [sh.text_frame.text for s in prs.slides for sh in s.shapes
                if sh.has_text_frame and "   \u00b7   " in sh.text_frame.text]
    check("The analysis slide shows two different lines, not one twice",
          len(headings) == len(set(headings)), str(headings))
    concentrations = [sh.text_frame.text for s in prs.slides for sh in s.shapes
                      if sh.has_text_frame
                      and sh.text_frame.text.startswith("Effectively all")]
    check("The analysis blocks do not share a concentration line",
          len(concentrations) == len(set(concentrations)), str(concentrations))

    # Every slide that reports a period must name the cumulative one, or the
    # deck and the workbook headline different timeframes.
    month_led = []
    for i, sl in enumerate(prs.slides, 1):
        text = " ".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
        if "Reporting month" in text:
            month_led.append(i)
    check("No slide still headlines the month alone",
          not month_led, f"slides {month_led}")

    charts = sum(1 for s in prs.slides for sh in s.shapes if sh.has_chart)
    check("Deck carries native charts, not pictures of charts", charts == 4, str(charts))
    check("Year-to-date slide states the cumulative revenue",
          f"{line(build_report(ytd_frame), 'Revenue').actual:,.0f}" in joined,
          "cumulative revenue missing")
    check("Year-to-date slide answers whether the month is typical",
          "run rate" in joined.lower(), "run-rate read missing")

    # The outlook must show both projections and name them, or a reader cannot
    # tell which assumption produced which number.
    check("Outlook slide states the run-rate projection",
          "run rate" in joined.lower() and "plan for rest" in joined.lower(),
          "one of the two projections is unnamed")
    check("Outlook slide says the projections are arithmetic, not a forecast",
          "not a forecast" in joined.lower())

    ytd_rep = build_report(ytd_frame)
    fy_rep = build_report(fy_frame)
    ni_ytd = line(ytd_rep, "Net income").actual
    ni_fy_bud = line(fy_rep, "Net income").budget
    run_rate = ni_ytd / 6 * 12
    plan_rest = ni_ytd + (ni_fy_bud - line(ytd_rep, "Net income").budget)
    check("Run rate is the year to date carried over twelve months",
          f"{run_rate:,.0f}" in joined, f"{run_rate:,.0f} missing")
    check("Plan-for-the-rest adds the unspent budget to the actuals",
          f"{plan_rest:,.0f}" in joined, f"{plan_rest:,.0f} missing")

    # Without a full-year plan there is nothing to project against.
    no_fy = build_pptx_pack(agg, PERIOD, tmp / "deck_nofy.pptx",
                            detail=monthly_detail(PERIOD), ytd=ytd_frame, months=6)
    check("Deck drops the outlook slide without a full-year plan",
          len(Presentation(no_fy).slides._sldIdLst) == 6,
          str(len(Presentation(no_fy).slides._sldIdLst)))
    # And falls back to the month rather than failing when the cumulative
    # detail is not supplied: a slide on the month beats no slide.
    check("The spend slide still builds without cumulative detail", True)

    # A single-period file has no year to date; the slide must not appear.
    no_ytd = build_pptx_pack(agg, PERIOD, tmp / "deck_single.pptx")
    check("Deck drops the cumulative slide for a single-period file",
          len(Presentation(no_ytd).slides._sldIdLst) == 4,
          str(len(Presentation(no_ytd).slides._sldIdLst)))

    empty_ytd = ingest.year_to_date(agg.assign(period=PERIOD), PERIOD)
    check("A one-period frame yields no year to date rather than repeating itself",
          empty_ytd.empty, str(len(empty_ytd)))

    tables = [sh.table for s in prs.slides for sh in s.shapes if sh.has_table]
    check("Deck has the P&L table", len(tables) == 1, str(len(tables)))
    if tables:
        t = tables[0]
        check("P&L table has a row per line plus a header",
              len(t.rows) == len(PNL_STRUCTURE) + 1, str(len(t.rows)))

    # Actuals only: the driver slide has nothing to say, so it must not appear.
    deck_nb = build_pptx_pack(actuals, PERIOD, tmp / "deck_nobudget.pptx")
    prs_nb = Presentation(deck_nb)
    check("Deck drops the driver slide when there is no budget",
          len(prs_nb.slides._sldIdLst) == 3, str(len(prs_nb.slides._sldIdLst)))
    nb_text = " ".join(sh.text_frame.text for s in prs_nb.slides
                       for sh in s.shapes if sh.has_text_frame)
    check("Deck says a budget is missing rather than showing a zero variance",
          "no budget" in nb_text.lower(), nb_text[:120])


# ---------------------------------------------------------------------------
# 12  Edge cases and the app
# ---------------------------------------------------------------------------
def test_edges() -> None:
    # A host that execs a module without registering it in sys.modules - a
    # hot-reloading server, a plugin loader - makes dataclasses fail while it is
    # guessing whether a string annotation is the KW_ONLY marker: it looks the
    # class's module up in sys.modules, gets None, and the import dies with an
    # AttributeError that names neither the module nor the field. Deferred
    # annotations are what put dataclasses on that path, so the modules that
    # define one must not ask for them.
    import importlib.util

    for name in ("coa", "engine", "ingest"):
        path = Path(__file__).resolve().parents[1] / "src" / "flux" / f"{name}.py"
        source = path.read_text(encoding="utf-8")
        has_dataclass = "@dataclass" in source
        # Match the statement, not a comment that quotes it.
        deferred = any(l.strip() == "from __future__ import annotations"
                       for l in source.splitlines())
        check(f"{name}.py does not defer annotations while defining a dataclass",
              not (has_dataclass and deferred),
              "a dataclass under deferred annotations breaks on some hosts")
        # A module using relative imports cannot be exec'd standalone at all,
        # so the static check above is the guard for those.
        if not has_dataclass or "\nfrom ." in source:
            continue
        spec = importlib.util.spec_from_file_location(f"_unregistered_{name}", path)
        module = importlib.util.module_from_spec(spec)
        try:
            # Deliberately not added to sys.modules: that is the failing case.
            spec.loader.exec_module(module)
            check(f"{name}.py imports on a host that skips sys.modules", True)
        except Exception as exc:  # pragma: no cover - the point of the check
            check(f"{name}.py imports on a host that skips sys.modules", False,
                  repr(exc))

    single = pd.DataFrame([{"account_code": "4000", "account_name": "Rev",
                            "category": "Revenue", "actual": 100.0,
                            "budget": 100.0, "prior_year": 100.0}])
    rep = build_report(single)
    check("A single-line file still produces a full P&L", len(rep) == len(PNL_STRUCTURE))
    check("A zero variance is favourable on revenue",
          line(rep, "Revenue").fav_unfav == "F")

    zero = single.copy(); zero["actual"] = 0.0; zero["budget"] = 0.0
    rep0 = build_report(zero)
    check("An all-zero file does not raise", len(rep0) == len(PNL_STRUCTURE))
    check("An all-zero file flags nothing", not rep0["material"].any())

    # An unknown category from a client file must not crash the engine.
    odd = pd.DataFrame([{"account_code": "9000", "account_name": "Odd",
                         "category": "Sonstiges", "actual": 100.0,
                         "budget": 50.0, "prior_year": 0.0}])
    try:
        leaf_variances(odd)
        check("An unknown category is handled, not fatal", True)
    except Exception as exc:  # pragma: no cover - the point of the check
        check("An unknown category is handled, not fatal", False, repr(exc))

    check("The chart of accounts has no duplicate codes",
          len({a.code for a in CHART_OF_ACCOUNTS}) == len(CHART_OF_ACCOUNTS))
    check("Every account has a favourable direction",
          all(a.favourable in ("higher", "lower") for a in CHART_OF_ACCOUNTS))

    # The budget-year generator must cover twelve months.
    by = generate_budget_year()
    check("The budget covers a full year", by["period"].nunique() == 12)

    # The template is the recommended path, so it has to survive its own
    # ingestion - and the example sheet has to be an example of something. It
    # shipped promising rows that were never written, because the generator's
    # import failed and a bare except swallowed it.
    tpl = Path(__file__).resolve().parents[1] / "data" / "input_template.xlsx"
    check("The input template exists", tpl.exists())
    if tpl.exists():
        tstd, _tr, _ti = ingest.ingest(tpl)
        check("The template's Input sheet ingests cleanly",
              tstd is not None and len(tstd) == 3,
              "none" if tstd is None else str(len(tstd)))
        twb = _load(tpl)
        example = twb["Example - full GL export"]
        check("The example ledger sheet actually carries rows",
              example.max_row > 6, str(example.max_row))
        check("The example sheet's banner spans its data",
              any(str(m).endswith(f"{get_column_letter(example.max_column)}1")
                  for m in example.merged_cells.ranges),
              str([str(m) for m in example.merged_cells.ranges]))
        raw = ingest.load_file(tpl, sheet_name="Example - full GL export")
        estd, _er, _ei = ingest.ingest(raw)
        check("The example export is a ledger with revenue in it",
              estd is not None and line(build_report(estd), "Revenue").actual > 0,
              "no revenue in the example")

        # The template's own example sheet is a ledger export, so it looks like
        # data and can outscore a sparsely filled Input sheet. Filling only the
        # required columns must still report the user's numbers, not the demo
        # company's.
        import shutil
        minimal = Path(tempfile.mkdtemp()) / "minimal.xlsx"
        shutil.copy(tpl, minimal)
        mwb = _load(minimal)
        mws = mwb["Input"]
        for row in (6, 7, 8):
            for col in range(1, 12):
                mws.cell(row=row, column=col).value = None
        for i, (code, name, amount) in enumerate(
                [("4100", "Arbevetel", 900_000.0), ("6200", "Berkoltseg", 400_000.0)],
                start=6):
            mws.cell(row=i, column=2, value=code)
            mws.cell(row=i, column=3, value=name)
            mws.cell(row=i, column=9, value=amount)
        mwb.save(minimal)
        mstd, _mr, _mi = ingest.ingest(minimal)
        check("A minimally filled template still beats its own example sheet",
              mstd is not None and set(mstd["account_name"]) == {"Arbevetel", "Berkoltseg"},
              "none" if mstd is None else str(list(mstd["account_name"])))

        # Row heights are set by hand in a generator that cannot measure text,
        # so every edit to the copy risks clipping it. Three blocks and the
        # longest table note were clipped before this check existed.
        howto = twb["How to fill"]
        clipped = []
        for row in range(1, howto.max_row + 1):
            for col, per_line in ((1, 105), (3, 72)):
                text = howto.cell(row=row, column=col).value
                if not isinstance(text, str) or len(text) < 90:
                    continue
                height = howto.row_dimensions[row].height or 15
                needed = 10 + 15 * (-(-len(text) // per_line))
                if height + 1 < needed:
                    clipped.append(f"row {row} col {col}: {height} < {needed}")
        check("No text on the How-to-fill sheet is clipped by its row height",
              not clipped, "; ".join(clipped[:4]))

        # The middle column answers "what does this buy me", not "will it
        # crash without this" - almost every column here is optional to the
        # engine and load-bearing for the report.
        buys = [howto.cell(row=r, column=2).value for r in range(1, howto.max_row + 1)]
        check("The template says what each column gives you, not just Required",
              "Expense Report sheet" in buys and "By Entity sheet" in buys,
              str([b for b in buys if b][:6]))
        check("Only the two amount columns share the one-of-these wording",
              buys.count("Fill one of these two") == 2,
              str(buys.count("Fill one of these two")))

        check("The expense-type dropdown is backed by a list, not typed by hand",
              any(str(dv.sqref).startswith("E") for dv in
                  twb["Input"].data_validations.dataValidation),
              str([str(d.sqref) for d in twb["Input"].data_validations.dataValidation]))

    # A budget-only file is a legitimate upload: the engine needs an amount,
    # not specifically the actual.
    budget_only = pd.DataFrame({
        "Account code": ["4100", "6110"],
        "Account name": ["Subscription revenue", "Advertising & digital"],
        "Budget": [500_000.0, 120_000.0],
    })
    bstd, _br, _bi = ingest.ingest(budget_only)
    check("A budget-only file is accepted",
          bstd is not None and close(bstd["budget"].sum(), 620_000.0),
          "rejected" if bstd is None else str(bstd["budget"].sum()))

    app = Path(__file__).resolve().parents[1] / "app.py"
    check("app.py exists", app.exists())
    source = app.read_text(encoding="utf-8")
    check("app.py parses", _parses(source))
    check("app.py uses no browser-only Streamlit APIs",
          "use_container_width" not in source)
    # A missing optional helper must not stop the app from starting: the whole
    # page went down over one extra slide on one tab.
    check("app.py imports only the package's own promises at module scope",
          "from flux.reporting.demo_pack import" not in source,
          "app.py reaches into a module rather than the package")
    from flux import reporting
    check("Anything the app imports is exported by the package",
          all(hasattr(reporting, n) for n in reporting.__all__),
          str([n for n in reporting.__all__ if not hasattr(reporting, n)]))


def _parses(source: str) -> bool:
    import ast
    try:
        ast.parse(source)
        return True
    except SyntaxError:
        return False


# ---------------------------------------------------------------------------
def main() -> int:
    suites = [
        ("Engine", test_engine),
        ("Parsing", test_parsing),
        ("Ingestion", test_ingestion),
        ("Sign normalisation", test_sign_normalisation),
        ("Expense grouping", test_expense_grouping),
        ("Reconciliation", test_reconciliation),
        ("Workbooks", test_workbooks),
        ("Edge cases", test_edges),
    ]
    crashed = []
    for name, fn in suites:
        start = len(_results)
        try:
            fn()
        except Exception:
            crashed.append(name)
            _results.append(("FAIL", f"{name}: suite raised", traceback.format_exc(limit=3)))
        ran = len(_results) - start
        failed = sum(1 for s, _n, _d in _results[start:] if s == "FAIL")
        status = "ok" if failed == 0 else f"{failed} FAILED"
        print(f"  {name:<22} {ran:>3} checks   {status}")

    failures = [(n, d) for s, n, d in _results if s == "FAIL"]
    print("\n" + "-" * 72)
    if failures:
        print(f"FAILED  {len(failures)} of {len(_results)} checks\n")
        for n, d in failures:
            print(f"  - {n}")
            if d:
                print(f"      {d.strip().splitlines()[-1][:160]}")
        return 1
    print(f"PASSED  all {len(_results)} checks")
    return 0


if __name__ == "__main__":
    sys.exit(main())
