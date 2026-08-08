"""
PowerPoint pack: the deck a controller sends upward after close.

The Excel pack is the working file - every figure a live formula over the input
sheet. This is the other half of the same job: five slides that state the
result, name what moved it, and stop. Same engine, same numbers, same palette.

Built with python-pptx rather than a Node deck library so the capability ships
inside the Python package with no extra runtime.

Run:  PYTHONPATH=src python -m flux.reporting.pptx_pack
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION,
                             XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION)
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..analysis import NO_CAUSE_NOTE
from ..commentary import generate_commentary, line_comments
from ..engine import (MaterialityRule, build_report, department_variances,
                      entity_variances, has_budget, leaf_variances)


# ---------------------------------------------------------------------------
# Palette - the same colours the Excel packs are rendered in.
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x24, 0x38)
DEEP = RGBColor(0x13, 0x1A, 0x29)
BRASS = RGBColor(0xC6, 0xA1, 0x5B)
IVORY = RGBColor(0xF5, 0xF1, 0xE8)
BAND = RGBColor(0xFA, 0xF8, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x23, 0x27, 0x2E)
MUTE = RGBColor(0x8A, 0x8F, 0x98)
PALE = RGBColor(0xD7, 0xDC, 0xE6)
SLATE = RGBColor(0x8A, 0xA0, 0xCC)
PALE_BAR = RGBColor(0xC8, 0xD0, 0xE0)
GREEN = RGBColor(0x2E, 0x6B, 0x3E)
RED = RGBColor(0xA6, 0x3A, 0x3A)
AMBER = RGBColor(0x8A, 0x5A, 0x00)
AMBER_FILL = RGBColor(0xFB, 0xEE, 0xDA)

# Calibri renders true to width in preview and ships with Office, so the
# slides measure the same here as they will on the reader's machine.
FONT = "Calibri"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)


# ---------------------------------------------------------------------------
# Small helpers
# ---------------------------------------------------------------------------
def _mag(v: float) -> str:
    """Money at the scale a slide is read from across a room."""
    a = abs(v)
    if a >= 1_000_000:
        return f"{v / 1_000_000:,.1f}m"
    if a >= 1_000:
        return f"{v / 1_000:,.0f}k"
    return f"{v:,.0f}"


def _signed(v: float) -> str:
    return ("+" if v >= 0 else "-") + _mag(abs(v)).lstrip("-")


def _pct(p) -> str:
    if p is None or (isinstance(p, float) and pd.isna(p)):
        return "n/m"
    return f"{p * 100:+.1f}%"


def _no_invert(series, idx: int) -> None:
    """Stop PowerPoint hollowing out a negative bar.

    A `c:dPt` carries its own `invertIfNegative`, and the schema default is true
    when the element is absent, so setting the flag on the series is not enough:
    the per-point default wins it back. python-pptx exposes no setter, and
    `Point._element` is the series rather than the point, so the data-point
    element is fetched explicitly and the flag inserted where CT_DPt requires
    it - straight after `c:idx`, before `c:spPr`.
    """
    from pptx.oxml.ns import qn

    dpt = series._element.get_or_add_dPt_for_point(idx)
    if dpt.find(qn("c:invertIfNegative")) is not None:
        return
    flag = dpt.makeelement(qn("c:invertIfNegative"), {"val": "0"})
    dpt.find(qn("c:idx")).addnext(flag)


def _label_points(plot, values=()) -> None:
    """Data labels, styled the same on every chart in the deck.

    Outside the bar, in ink: a category can be a rounding error against the
    others, and a label placed inside that sliver would be white on white.

    The scale follows the numbers. Full digits on a chart of millions run three
    adjacent labels into each other, so past a million they read in millions -
    chosen by magnitude rather than per chart, so no two slides format the same
    figure differently.
    """
    plot.has_data_labels = True
    labels = plot.data_labels
    biggest = max((abs(float(v)) for v in values), default=0.0)
    labels.number_format = ('#,##0.0,,"m";(#,##0.0,,"m")' if biggest >= 1_000_000
                            else '#,##0;(#,##0)')
    labels.number_format_is_linked = False
    labels.font.size = Pt(9)
    labels.font.name = FONT
    labels.font.bold = True
    labels.font.color.rgb = INK
    labels.position = XL_LABEL_POSITION.OUTSIDE_END


def _pad_value_axis(chart, values, pad=0.28) -> None:
    """Leave room for the outside labels, on bounds a reader can hold.

    Padding a raw minimum gives an axis that starts at -101,752 and ticks in
    steps to match, which is harder to read than the numbers it is scaling. The
    padded bounds are rounded out to a round step first.
    """
    import math

    lo, hi = min(list(values) + [0]), max(list(values) + [0])
    span = (hi - lo) or 1.0
    lo = lo - span * pad if lo < 0 else 0
    hi = hi + span * pad if hi > 0 else 0
    step = 10 ** max(0, int(math.floor(math.log10(span))) - 1) * 2
    chart.value_axis.minimum_scale = math.floor(lo / step) * step if lo else 0
    chart.value_axis.maximum_scale = math.ceil(hi / step) * step if hi else 0


def _axis_style(chart) -> None:
    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(10)
        axis.tick_labels.font.name = FONT
        axis.tick_labels.font.color.rgb = INK
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE4, 0xE7, 0xEC)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.75)
    chart.value_axis.tick_labels.number_format = '#,##0;(#,##0)'
    chart.value_axis.tick_labels.number_format_is_linked = False


def _blank(prs: Presentation):
    """A slide with nothing on it: every element here is placed explicitly."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(shape, colour, line=None, width=Pt(1)):
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = width
    shape.shadow.inherit = False


def _text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          spacing=None):
    """Place a text box. `runs` is a list of (text, size, bold, colour) tuples;
    a None entry starts a new paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    for item in runs:
        if item is None:
            para = tf.add_paragraph()
            para.alignment = align
            if spacing is not None:
                para.space_before = spacing
            continue
        txt, size, bold, colour = item
        run = para.add_run()
        run.text = txt
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
    return box


def _mark(slide, x, y, scale=1.0, plan=SLATE, actual=BRASS):
    """The Flux mark, drawn as shapes.

    Three rectangles: the plan bar, the shorter actual bar, and the rule showing
    where plan sat. Native shapes rather than an image, so it stays sharp at any
    size and the package needs no renderer to ship a logo.
    """
    u = Inches(0.028) * scale                       # the mark's grid unit
    def r(gx, gy, gw, gh, colour):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Emu(int(u * gx)), y + Emu(int(u * gy)),
                                   Emu(int(u * gw)), Emu(int(u * gh)))
        s.adjustments[0] = 0.18
        _fill(s, colour)
        return s
    r(0, 0, 15, 30, plan)
    r(19, 12, 15, 18, actual)
    r(19, 0, 15, 4, actual)


def _slide_header(slide, eyebrow, title, meta=""):
    """Light content slide: small brass eyebrow, then the headline."""
    _text(slide, MARGIN, Inches(0.46), Inches(9.0), Inches(0.3),
          [(eyebrow.upper(), 11, True, BRASS)])
    _text(slide, MARGIN, Inches(0.74), Inches(9.6), Inches(0.6),
          [(title, 30, True, NAVY)])
    if meta:
        # Below the mark, not beside it: at this size the mark is only 0.4" wide
        # and the two collided on every content slide.
        _text(slide, W - MARGIN - Inches(4.2), Inches(1.02), Inches(4.2), Inches(0.3),
              [(meta, 12, False, MUTE)], align=PP_ALIGN.RIGHT)
    _mark(slide, W - MARGIN - Inches(0.42), Inches(0.46), scale=0.36,
          plan=RGBColor(0xC8, 0xD0, 0xE0), actual=BRASS)


def _dark_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _fill(bg, DEEP)
    return bg


# ---------------------------------------------------------------------------
# Slide 1 - cover
# ---------------------------------------------------------------------------
def _cover(prs, entity, period, ni, budgeted, ytd_ni=None):
    s = _blank(prs)
    _dark_bg(s)
    _mark(s, MARGIN, Inches(1.62), scale=1.45)

    word_x = MARGIN + Inches(1.85)
    _text(s, word_x, Inches(1.60), Inches(7.5), Inches(0.9),
          [("FLUX", 52, True, BRASS)])
    _text(s, word_x, Inches(2.44), Inches(9.0), Inches(0.4),
          [("Management reporting pack", 19, False, PALE)])

    _text(s, MARGIN, Inches(3.75), Inches(8.0), Inches(0.5),
          [(entity, 30, True, WHITE)])
    meta = (f"YTD through {period}, with the month  ·  currency EUR"
            if ytd_ni is not None else f"Reporting month {period}  ·  currency EUR")
    _text(s, MARGIN, Inches(4.35), Inches(8.0), Inches(0.4),
          [(meta, 15, False, MUTE)])

    # The headline number belongs on the cover: it is what the reader opened for.
    # The headline is the year to date, as it is everywhere else in the pack:
    # the cover used to show the month, so the first number a reader saw was the
    # one the rest of the deck treats as the secondary reading.
    head = ytd_ni if ytd_ni is not None else ni
    label = "NET INCOME · YTD" if ytd_ni is not None else "NET INCOME"
    tall = Inches(1.85) if ytd_ni is not None else Inches(1.55)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              W - MARGIN - Inches(4.0), Inches(3.55),
                              Inches(4.0), tall)
    card.adjustments[0] = 0.09
    _fill(card, NAVY, line=BRASS, width=Pt(1.25))
    _text(s, W - MARGIN - Inches(3.7), Inches(3.78), Inches(3.4), Inches(0.28),
          [(label, 11, True, BRASS)])
    _text(s, W - MARGIN - Inches(3.7), Inches(4.06), Inches(3.4), Inches(0.55),
          [(f"{head['actual']:,.0f} EUR", 30, True, WHITE)])
    sub = (f"{_signed(head['var_bud'])} ({_pct(head['var_bud_pct'])}) vs budget"
           if budgeted else "no budget supplied")
    _text(s, W - MARGIN - Inches(3.7), Inches(4.66), Inches(3.4), Inches(0.3),
          [(sub, 12, False, RED if budgeted and head["var_bud"] < 0 else PALE)])
    if ytd_ni is not None:
        month = (f"Month {ni['actual']:,.0f}"
                 + (f"  ·  {_pct(ni['var_bud_pct'])}" if budgeted else ""))
        _text(s, W - MARGIN - Inches(3.7), Inches(4.99), Inches(3.4), Inches(0.3),
              [(month, 11, False, MUTE)])

    _text(s, MARGIN, H - Inches(0.95), Inches(11.0), Inches(0.3),
          [("Generated by Flux from the general ledger. Every figure ties to the "
            "Excel pack issued with this deck.", 11, False, MUTE)])
    return s


# ---------------------------------------------------------------------------
# Slide 2 - the result
# ---------------------------------------------------------------------------
def _kpi_card(s, x, y, w, label, value, delta, favourable, budgeted,
              second=""):
    """A KPI card: the headline figure, its delta, and the month beneath.

    The workbook headlines the year to date and carries the month as the second
    reading, because a single month is noisy and cumulative performance against
    the annual plan is what a management cover is read for. The deck follows it,
    so the two cannot show the same company two ways.
    """
    height = Inches(1.78) if second else Inches(1.5)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, height)
    card.adjustments[0] = 0.09
    _fill(card, IVORY, line=BRASS, width=Pt(1))
    _text(s, x + Inches(0.26), y + Inches(0.2), w - Inches(0.5), Inches(0.25),
          [(label.upper(), 11, True, MUTE)])
    _text(s, x + Inches(0.26), y + Inches(0.48), w - Inches(0.5), Inches(0.5),
          [(f"{value:,.0f}", 28, True, NAVY)])
    if budgeted:
        _text(s, x + Inches(0.26), y + Inches(1.06), w - Inches(0.5), Inches(0.28),
              [(delta, 12, True, GREEN if favourable else RED)])
    if second:
        _text(s, x + Inches(0.26), y + Inches(1.38), w - Inches(0.5), Inches(0.28),
              [(second, 11, False, MUTE)])


def _result(prs, report, gl, period, budgeted, ytd_report=None):
    s = _blank(prs)
    headline = ytd_report if ytd_report is not None else report
    subtitle = ("Where the year has got to" if ytd_report is not None
                else "Where the month landed")
    meta = (f"YTD through {period}, with the month  ·  EUR" if ytd_report is not None
            else f"Reporting month {period}  ·  EUR")
    _slide_header(s, "The result", subtitle, meta)

    def line(frame, label):
        return frame[frame.line == label].iloc[0]

    labels = ("Revenue", "Operating income (EBIT)", "Net income")
    gap = Inches(0.34)
    cw = (W - 2 * MARGIN - 2 * gap) / 3
    for i, label in enumerate(labels):
        head = line(headline, label)
        second = ""
        if ytd_report is not None:
            m = line(report, label)
            second = (f"Month {m.actual:,.0f}"
                      + (f"  ·  {_pct(m.var_bud_pct)}" if budgeted else ""))
        _kpi_card(s, MARGIN + i * (cw + gap), Inches(1.62), cw,
                  (head.line + " · YTD") if ytd_report is not None else head.line,
                  head.actual,
                  f"{_signed(head.var_bud)} ({_pct(head.var_bud_pct)}) vs budget"
                  if budgeted else "", head.fav_unfav == "F", budgeted, second)

    # The narrative's first two paragraphs: the bottom line, then revenue and margin.
    text = generate_commentary(report, gl).split("\n\n")
    body = s.shapes.add_textbox(MARGIN, Inches(3.72), W - 2 * MARGIN, Inches(3.1))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, para_text in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = para_text
        run.font.name = FONT
        run.font.size = Pt(15)
        run.font.color.rgb = INK
    s.notes_slide.notes_text_frame.text = (
        "Figures are the reporting month against budget. The commentary is "
        "generated from the account-level variances, not written by hand."
    )
    return s


# ---------------------------------------------------------------------------
# Slide 3 - the P&L
# ---------------------------------------------------------------------------
def _pnl_table(prs, report, comments, period, budgeted, ytd_report=None):
    s = _blank(prs)
    both = ytd_report is not None and budgeted
    meta = (f"YTD through {period}, with the month  ·  EUR" if both
            else f"Reporting month {period}  ·  EUR")
    _slide_header(s, "Profit and loss", "The management P&L", meta)

    # The workbook reports both timeframes on every line; a deck that showed
    # only the month made the same statement two different ways.
    if both:
        heads = ["Line", "YTD Act", "YTD Var", "Var %", "Month Act",
                 "Month Var", "F/U", ""]
    elif budgeted:
        heads = ["Line", "Actual", "Budget", "Variance", "Var %", "F/U", ""]
    else:
        heads = ["Line", "Actual", "Prior year", "Var (PY)", "", "", ""]
    rows = len(report) + 1
    # The money columns hold at most nine characters, the commentary a sentence,
    # so width goes where the text actually is.
    widths = ([Inches(2.3), Inches(1.2), Inches(1.15), Inches(0.85),
               Inches(1.2), Inches(1.15), Inches(0.5), Inches(3.74)] if both
              else [Inches(2.45), Inches(1.25), Inches(1.25), Inches(1.3),
                    Inches(0.95), Inches(0.55), Inches(4.34)])

    tbl_shape = s.shapes.add_table(rows, len(heads), MARGIN, Inches(1.66),
                                   sum(widths, Inches(0)), Inches(0.48) * rows)
    tbl = tbl_shape.table
    tbl.first_row = False
    for i, wd in enumerate(widths):
        tbl.columns[i].width = wd
    tbl.rows[0].height = Inches(0.4)

    def cell(r, c, txt, *, size=12, bold=False, colour=INK, align=PP_ALIGN.RIGHT,
             fill=None):
        cl = tbl.cell(r, c)
        cl.text = ""
        cl.margin_left = cl.margin_right = Inches(0.09)
        cl.margin_top = cl.margin_bottom = Inches(0.02)
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fill is not None:
            cl.fill.solid(); cl.fill.fore_color.rgb = fill
        else:
            cl.fill.background()
        p = cl.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour

    for c, h in enumerate(heads):
        cell(0, c, h, size=11, bold=True, colour=WHITE, fill=NAVY,
             align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)

    subs = {"Gross profit", "Operating income (EBIT)", "Net income"}
    for i, r in enumerate(report.itertuples(), start=1):
        sub = r.line in subs
        bg = IVORY if sub else (BAND if i % 2 else WHITE)
        tbl.rows[i].height = Inches(0.48)
        cell(i, 0, r.line, bold=sub, align=PP_ALIGN.LEFT, fill=bg)
        if both:
            y = ytd_report[ytd_report.line == r.line].iloc[0]
            cell(i, 1, f"{y.actual:,.0f}", bold=sub, fill=bg)
            cell(i, 2, f"{y.var_bud:,.0f}", bold=sub, fill=bg,
                 colour=GREEN if y.fav_unfav == "F" else RED)
            cell(i, 3, _pct(y.var_bud_pct), bold=sub, fill=bg)
            cell(i, 4, f"{r.actual:,.0f}", bold=sub, fill=bg, colour=MUTE)
            cell(i, 5, f"{r.var_bud:,.0f}", bold=sub, fill=bg, colour=MUTE)
            cell(i, 6, r.fav_unfav, bold=True, align=PP_ALIGN.CENTER, fill=bg,
                 colour=GREEN if r.fav_unfav == "F" else RED)
            note = comments.get(r.line, "")
            cell(i, 7, note if len(note) <= 118 else note[:115].rstrip(" ,;") + "...",
                 size=9, colour=MUTE, align=PP_ALIGN.LEFT, fill=bg)
            continue
        cell(i, 1, f"{r.actual:,.0f}", bold=sub, fill=bg)
        if budgeted:
            cell(i, 2, f"{r.budget:,.0f}", bold=sub, fill=bg)
            cell(i, 3, f"{r.var_bud:,.0f}", bold=sub, fill=bg,
                 colour=GREEN if r.fav_unfav == "F" else RED)
            cell(i, 4, _pct(r.var_bud_pct), bold=sub, fill=bg)
            cell(i, 5, r.fav_unfav, bold=True, align=PP_ALIGN.CENTER, fill=bg,
                 colour=GREEN if r.fav_unfav == "F" else RED)
        else:
            cell(i, 2, f"{r.prior_year:,.0f}", bold=sub, fill=bg)
            cell(i, 3, f"{r.var_py:,.0f}", bold=sub, fill=bg)
            cell(i, 4, "", fill=bg); cell(i, 5, "", fill=bg)
        note = comments.get(r.line, "")
        cell(i, 6, note if len(note) <= 132 else note[:129].rstrip(" ,;") + "...",
             size=10, colour=MUTE, align=PP_ALIGN.LEFT, fill=bg)

    if budgeted:
        tail = (" The year to date leads and the month is muted beside it, as in "
                "the workbook; F/U judges the month." if both else "")
        _text(s, MARGIN, H - Inches(0.82), Inches(11.5), Inches(0.3),
              [("Favourable / unfavourable follows the account type: revenue and "
                "profit lines are higher-is-better, cost lines lower-is-better."
                + tail, 10, False, MUTE)])
    return s


# ---------------------------------------------------------------------------
# Slide 4 - what moved it
# ---------------------------------------------------------------------------
def _year_to_date(prs, ytd_report, month_report, months, period, budgeted):
    """The cumulative view, and whether the month is typical of it.

    A repeat of the monthly table under a different heading would add pages
    without adding an answer. The question a cumulative slide exists to settle
    is whether the month is a blip or the run rate, so the chart puts the
    month's variance beside the average month of the year so far.
    """
    s = _blank(prs)
    _slide_header(s, "Year to date", f"{months} months in",
                  f"Through {period}  ·  EUR")

    def line(report, label):
        return report[report.line == label].iloc[0]

    # No KPI cards here: the result slide already carries the same three
    # year-to-date figures, and repeating them made two slides that looked
    # identical and answered the same question twice. This slide exists for the
    # comparison beneath, which is the one thing the cards cannot show.
    _text(s, MARGIN, Inches(1.62), W - 2 * MARGIN, Inches(0.45),
          [("How the month compares with the average month so far, by P&L line. "
            "The year-to-date figures themselves are on the result slide.",
            13, False, MUTE)])

    cats = ["Revenue", "Cost of goods sold", "Operating expenses", "Other expenses"]
    this_month = [float(line(month_report, c).var_bud) for c in cats]
    average = [float(line(ytd_report, c).var_bud) / months for c in cats]

    data = CategoryChartData()
    data.categories = cats
    data.add_series("Average month, year to date", tuple(average))
    data.add_series(f"This month ({period})", tuple(this_month))

    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, Inches(2.25),
                            Inches(8.1), Inches(4.4), data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT
    chart.plots[0].gap_width = 70
    _label_points(chart.plots[0], list(average) + list(this_month))
    for series, colour in zip(chart.series, (SLATE, BRASS)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colour
        series.invert_if_negative = False
    _axis_style(chart)
    _pad_value_axis(chart, list(average) + list(this_month))

    # Name the read rather than leaving it to be inferred from two bar heights.
    ni_m = line(month_report, "Net income")
    ni_y = line(ytd_report, "Net income")
    run_rate = ni_y.var_bud / months
    worse = abs(ni_m.var_bud) > abs(run_rate) * 1.15
    better = abs(ni_m.var_bud) < abs(run_rate) * 0.85
    verdict = ("worse than the run rate" if worse
               else "better than the run rate" if better
               else "in line with the run rate")

    # Which category departs furthest from its own run rate: the chart shows the
    # gaps, this names the one worth asking about.
    spreads = []
    for c, m_var, avg in zip(cats, this_month, average):
        if abs(avg) >= 1_000:
            spreads.append((abs(m_var / avg), c, m_var, avg))
    if spreads:
        ratio, cat, m_var, avg = max(spreads)
        outlier = (f"{cat}: {_mag(abs(m_var))} this month against an average "
                   f"month of {_mag(abs(avg))} — {ratio:.1f}x.")
    else:
        outlier = "No category departs materially from its own run rate."

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               MARGIN + Inches(8.45), Inches(2.25),
                               Inches(3.2), Inches(4.4))
    panel.adjustments[0] = 0.07
    _fill(panel, IVORY, line=BRASS, width=Pt(1))
    runs = [("IS THE MONTH TYPICAL?", 11, True, BRASS), None,
            (verdict.capitalize(), 17, True, NAVY), None,
            (f"Net income is {_mag(abs(ni_y.var_bud))} behind plan over {months} "
             f"months, an average of {_mag(abs(run_rate))} a month. This month was "
             f"{_mag(abs(ni_m.var_bud))}.", 12, False, INK), None,
            (f"Year to date {_signed(ni_y.var_bud)}", 12, True,
             GREEN if ni_y.fav_unfav == "F" else RED), None,
            ("THE LINE THAT DIVERGED MOST", 11, True, BRASS), None,
            (outlier, 12, False, INK)]
    _text(s, MARGIN + Inches(8.7), Inches(2.52), Inches(2.7), Inches(3.9), runs,
          spacing=Pt(9))

    s.notes_slide.notes_text_frame.text = (
        "The cumulative variance divided by months elapsed gives the average "
        "month. Comparing this month against it separates a one-off from a trend."
    )
    return s


def _outlook(prs, ytd_report, fy_report, months, period):
    """Where the year lands, on two stated assumptions.

    A single projected number invites the question it hides: projected how? So
    both are shown. Run rate carries the year so far forward; plan-for-the-rest
    assumes the remaining months hit budget. The outcome usually sits between
    them, and the gap between the two is itself the message - it is the size of
    what the remaining months have to fix.
    """
    s = _blank(prs)
    _slide_header(s, "Outlook", "Where the year lands",
                  f"FY {period[:4]}  ·  EUR")

    def line(report, label):
        return report[report.line == label].iloc[0]

    cats = ["Revenue", "Gross profit", "Operating income (EBIT)", "Net income"]
    fy_budget = [float(line(fy_report, c).budget) for c in cats]
    run_rate = [float(line(ytd_report, c).actual) / months * 12 for c in cats]
    plan_rest = [float(line(ytd_report, c).actual)
                 + (b - float(line(ytd_report, c).budget))
                 for c, b in zip(cats, fy_budget)]

    data = CategoryChartData()
    data.categories = ["Revenue", "Gross profit", "EBIT", "Net income"]
    data.add_series("FY budget", tuple(fy_budget))
    data.add_series("If the rest goes to plan", tuple(plan_rest))
    data.add_series("If the year runs as so far", tuple(run_rate))

    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, Inches(1.66),
                            Inches(8.1), Inches(4.6), data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT
    chart.plots[0].gap_width = 80
    _label_points(chart.plots[0], fy_budget + run_rate + plan_rest)
    for series, colour in zip(chart.series, (PALE_BAR, SLATE, BRASS)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colour
        series.invert_if_negative = False
    _axis_style(chart)
    _pad_value_axis(chart, fy_budget + run_rate + plan_rest, pad=0.16)

    ni_bud = line(fy_report, "Net income").budget
    ni_run = run_rate[-1]
    ni_plan = plan_rest[-1]
    gap = ni_plan - ni_run

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               MARGIN + Inches(8.45), Inches(1.66),
                               Inches(3.2), Inches(4.6))
    panel.adjustments[0] = 0.05
    _fill(panel, IVORY, line=BRASS, width=Pt(1))
    runs = [("NET INCOME, FULL YEAR", 11, True, BRASS), None,
            (f"Budget {ni_bud:,.0f}", 13, True, NAVY), None,
            (f"Run rate {ni_run:,.0f}", 12, False,
             GREEN if ni_run >= ni_bud else RED), None,
            (f"Plan for rest {ni_plan:,.0f}", 12, False,
             GREEN if ni_plan >= ni_bud else RED), None,
            ("THE GAP BETWEEN THEM", 11, True, BRASS), None,
            (f"{_mag(abs(gap))}", 17, True, NAVY), None,
            (f"That is what the remaining {12 - months} months have to make up "
             f"for the year to land on the better of the two paths.",
             11, False, INK), None,
            ("Both are arithmetic, not a forecast: no seasonality, no pipeline, "
             "no management action.", 10, False, MUTE)]
    _text(s, MARGIN + Inches(8.7), Inches(1.92), Inches(2.7), Inches(4.1), runs,
          spacing=Pt(8))

    s.notes_slide.notes_text_frame.text = (
        "Run rate = year to date divided by months elapsed, times twelve. "
        "Plan for the rest = actual to date plus the unspent part of the budget. "
        "Neither is a forecast; they bracket the outcome under a stated assumption."
    )
    return s


# ---------------------------------------------------------------------------
# Slide - what the numbers point at
# ---------------------------------------------------------------------------
def _analysis(prs, blocks, period):
    """The same analysis the workbook carries under each sheet.

    One slide, the two lines that carry a flag, and for each the shape of the
    movement and the question it makes worth asking. Never a cause: the ledger
    does not record one, and a deck is the last place to invent it.
    """
    s = _blank(prs)
    _slide_header(s, "Analysis", "What the numbers point at",
                  f"YTD through {period}, with the month  \u00b7  EUR")

    # EBIT and net income produce the same findings - every category's movement
    # lands on both - so showing them one under the other filled the slide with
    # one thought written twice. Blocks whose findings match are dropped.
    # Keyed on everything except the full-year line: EBIT and net income differ
    # only in the size of the plan they are measured against, which is not a
    # second thought. Matching on the whole block would have kept both.
    seen, distinct = set(), []
    for label, flag, items in blocks:
        key = tuple(f.text for f in items if f.heading != "Full year")
        if key in seen:
            continue
        seen.add(key)
        distinct.append((label, flag, items))

    y = Inches(1.55)
    for label, flag, items in distinct[:2]:
        head = f"{label}   \u00b7   {flag}" if flag else label
        _text(s, MARGIN, y, W - 2 * MARGIN, Inches(0.3), [(head, 15, True, NAVY)])
        y += Inches(0.36)
        for finding in items:
            _text(s, MARGIN + Inches(0.1), y, Inches(1.35), Inches(0.3),
                  [(finding.heading, 11, True, MUTE)])
            box = s.shapes.add_textbox(MARGIN + Inches(1.5), y - Inches(0.04),
                                       W - 2 * MARGIN - Inches(1.5), Inches(0.7))
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = 0
            run = tf.paragraphs[0].add_run()
            run.text = finding.text
            run.font.name = FONT
            run.font.size = Pt(11.5)
            run.font.color.rgb = INK
            y += Inches(0.34) if len(finding.text) < 115 else Inches(0.56)
        y += Inches(0.26)

    _text(s, MARGIN, H - Inches(1.0), W - 2 * MARGIN, Inches(0.6),
          [(NO_CAUSE_NOTE.replace("on this sheet", "in this pack"), 9, False, MUTE)])
    s.notes_slide.notes_text_frame.text = (
        "Derived from the figures in the pack: concentration, persistence, the "
        "full-year run rate and the question the shape of the variance implies. "
        "No causal claim is made anywhere in this deck."
    )
    return s


# ---------------------------------------------------------------------------
# Slide - what moved it
# ---------------------------------------------------------------------------
def _drivers(prs, gl, materiality, period, budgeted, ytd=None):
    """The accounts that moved the result, cumulatively where the data allows.

    The workbook's Drivers sheet reports the year to date beside the month, so a
    deck showing only the month named a different set of accounts from the file
    it ships with.
    """
    s = _blank(prs)
    frame = ytd if ytd is not None else gl
    meta = (f"YTD through {period}  \u00b7  EUR" if ytd is not None
            else f"Reporting month {period}  \u00b7  EUR")
    _slide_header(s, "Drivers", "What moved the result", meta)

    leaves = leaf_variances(frame, materiality)
    material = leaves[leaves["material"]].copy()
    if material.empty:
        _text(s, MARGIN, Inches(2.4), Inches(9.0), Inches(0.6),
              [("No account cleared both materiality floors this month.",
                18, False, INK)])
        return s

    material["mag"] = material["var_bud"].abs()
    material = material.sort_values("mag", ascending=False).head(8)
    # Charted bottom-up so the largest mover sits at the top of the bar chart.
    plot = material.iloc[::-1]

    data = CategoryChartData()
    data.categories = [n if len(n) < 34 else n[:31] + "..." for n in plot["account_name"]]
    data.add_series("Variance vs budget", tuple(plot["var_bud"]))

    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, MARGIN, Inches(1.62),
                            Inches(8.1), Inches(4.6), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot_area = chart.plots[0]
    plot_area.gap_width = 55
    _label_points(plot_area, list(plot["var_bud"]))

    # PowerPoint's default for a bar series is to invert the fill on negative
    # values, and OOXML treats the flag as true when it is absent. LibreOffice
    # ignores it, so the shortfall bars looked correct in preview and rendered
    # hollow in PowerPoint itself.
    series = chart.series[0]
    series.invert_if_negative = False

    # One colour per bar: a cost overrun and a revenue shortfall are both bad
    # news but have opposite signs, so sign alone cannot carry the meaning.
    for idx, fu in enumerate(plot["fav_unfav"]):
        point = series.points[idx]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = GREEN if fu == "F" else RED
        _no_invert(series, idx)

    _axis_style(chart)
    _pad_value_axis(chart, list(plot["var_bud"]))

    # The side panel says what the chart cannot: why these lines and no others.
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               MARGIN + Inches(8.45), Inches(1.62),
                               Inches(3.2), Inches(4.6))
    panel.adjustments[0] = 0.05
    _fill(panel, IVORY, line=BRASS, width=Pt(1))

    # Read the drivers by their effect on profit, not by the sign of the raw
    # variance. Revenue below plan and cost above plan both cut profit but carry
    # opposite signs, so adding the figures as printed would look like adding
    # unlike things. Each line's profit impact is the variance signed against
    # the result, and those do add up.
    rev_short = -material.loc[material.category == "Revenue", "var_bud"].sum()
    cost_over = material.loc[material.category != "Revenue", "var_bud"].sum()
    impact = -(rev_short + cost_over)

    pieces = []
    if rev_short:
        pieces.append(f"revenue {_mag(abs(rev_short))} "
                      f"{'short of' if rev_short > 0 else 'ahead of'} plan")
    if cost_over:
        pieces.append(f"costs {_mag(abs(cost_over))} "
                      f"{'over' if cost_over > 0 else 'under'} plan")
    total = ("Effect on profit: " + " and ".join(pieces) + "."
             if pieces else "No net effect on profit.")
    split = f"{_signed(impact)} on net income."

    runs = [("MATERIALITY", 11, True, BRASS), None,
            (f"{materiality.abs_threshold:,.0f} EUR and "
             f"{materiality.pct_threshold:.0%}", 17, True, NAVY), None,
            ("An account is flagged only when it clears both floors, so a large "
             "percentage on a small base does not crowd out the real movers.",
             11, False, INK), None,
            (total, 11, False, INK), None,
            (split, 13, True, GREEN if impact >= 0 else RED), None,
            ("LARGEST SINGLE MOVER", 11, True, BRASS), None,
            (str(material.iloc[0]["account_name"]), 13, True, NAVY), None,
            (f"{_signed(material.iloc[0]['var_bud'])}  "
             f"({_pct(material.iloc[0]['var_bud_pct'])} vs budget)", 12, False,
             GREEN if material.iloc[0]["fav_unfav"] == "F" else RED)]
    _text(s, MARGIN + Inches(8.7), Inches(1.88), Inches(2.7), Inches(4.1), runs,
          spacing=Pt(9))
    return s


# ---------------------------------------------------------------------------
# Slide 5 - where the money went
# ---------------------------------------------------------------------------
def _spend(prs, detail, period, materiality, ytd_detail=None):
    """Departmental spend and the entity consolidation.

    Reported on the year to date when it is available, with the month named
    underneath each entity - the same order of reading as every other page.
    """
    s = _blank(prs)
    frame = ytd_detail if ytd_detail is not None else detail
    if ytd_detail is not None and "prior_year" not in frame.columns:
        # The cumulative frame is built for this pack and carries no prior year;
        # the engine expects the column, so it is supplied as zero rather than
        # letting a missing comparison take the slide down.
        frame = frame.assign(prior_year=0.0)
    meta = (f"YTD through {period}, with the month  \u00b7  EUR"
            if ytd_detail is not None else f"Reporting month {period}  \u00b7  EUR")
    _slide_header(s, "Spend", "Where the money went", meta)

    dept = department_variances(frame, materiality).sort_values("actual", ascending=False)
    plot = dept.iloc[::-1]

    data = CategoryChartData()
    data.categories = list(plot["department"])
    # Budget first so the clustered bars put actual on top, where the eye lands.
    data.add_series("Budget", tuple(plot["budget"]))
    data.add_series("Actual", tuple(plot["actual"]))

    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, MARGIN, Inches(1.62),
                            Inches(7.4), Inches(4.5), data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = FONT
    chart.plots[0].gap_width = 60
    _label_points(chart.plots[0], list(plot["actual"]) + list(plot["budget"]))
    for series, colour in zip(chart.series, (SLATE, BRASS)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colour
        series.invert_if_negative = False
    _axis_style(chart)
    _pad_value_axis(chart, list(plot["actual"]) + list(plot["budget"]), pad=0.18)

    # Entity consolidation beside it: the same spend, cut the other way.
    ent = entity_variances(frame)
    ent_month = entity_variances(detail) if ytd_detail is not None else None
    x0 = MARGIN + Inches(7.75)
    _text(s, x0, Inches(1.66), Inches(3.9), Inches(0.3),
          [("BY LEGAL ENTITY", 11, True, BRASS)])
    y = Inches(2.06)
    for r in ent.itertuples():
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x0, y,
                                  Inches(3.9), Inches(1.02))
        card.adjustments[0] = 0.12
        _fill(card, BAND, line=RGBColor(0xE4, 0xE7, 0xEC), width=Pt(1))
        _text(s, x0 + Inches(0.22), y + Inches(0.14), Inches(2.2), Inches(0.28),
              [(r.entity, 14, True, NAVY)])
        month_note = ""
        if ent_month is not None:
            row = ent_month[ent_month["entity"] == r.entity]
            if len(row):
                month_note = f"month {row.iloc[0]['net_actual']:,.0f}"
        _text(s, x0 + Inches(0.22), y + Inches(0.46), Inches(2.4), Inches(0.42),
              [(f"Net income {r.net_actual:,.0f}"
                + (f"   ({month_note})" if month_note else ""), 12, False, INK), None,
               (f"against a plan of {r.net_budget:,.0f}", 10, False, MUTE)])
        # Net income is higher-is-better, so earning less than plan is the
        # unfavourable direction even though the variance reads negative. The
        # label spells out the comparison, which the bare number did not.
        fu_colour = GREEN if r.fav_unfav == "F" else RED
        _text(s, x0 + Inches(2.42), y + Inches(0.3), Inches(1.28), Inches(0.5),
              [(_signed(r.var_bud), 16, True, fu_colour), None,
               ("vs budget", 9, False, MUTE)], align=PP_ALIGN.RIGHT)
        y = y + Inches(1.18)

    _text(s, MARGIN, H - Inches(0.82), Inches(11.5), Inches(0.3),
          [("Departmental figures are spend only: revenue is excluded, so the "
            "totals tie to the cost lines of the P&L, not to net income.",
            10, False, MUTE)])
    return s


# ---------------------------------------------------------------------------
def build_pptx_pack(gl: pd.DataFrame,
                    period: str,
                    out_path: str | Path,
                    detail: pd.DataFrame | None = None,
                    ytd: pd.DataFrame | None = None,
                    fy_budget: pd.DataFrame | None = None,
                    months: int | None = None,
                    entity: str = "Demo Company Ltd",
                    materiality: MaterialityRule | None = None,
                    budgeted: bool | None = None,
                    analysis_blocks: list | None = None,
                    ytd_detail: pd.DataFrame | None = None) -> Path:
    """Build the management deck.

    Parameters
    ----------
    gl : account-level frame, the same shape the Excel pack is built from.
    period : reporting period label, e.g. "2025-06".
    detail : optional transaction-level frame carrying department and entity.
        Without it the spend slide is omitted rather than faked.
    ytd : optional account-level frame cumulative to `period`. Without it the
        cumulative and outlook slides are omitted; a single-period file has no
        year to date, and repeating the month under that heading would be a lie.
    fy_budget : optional account-level frame carrying the full-year plan. Needed
        for the outlook slide, which compares two projections against it.
    months : periods elapsed, used for the average month. Derived as needed.
    budgeted : whether a plan exists; detected from the data when omitted.
    ytd_detail : the dimensional detail cumulatively, at the same grain as
        `detail`. Without it the spend slide falls back to the month.
    analysis_blocks : optional findings from `flux.analysis`, the same ones the
        workbook prints under each sheet. Without them the analysis slide is
        omitted rather than filled with a restatement of the table.
    """
    materiality = materiality or MaterialityRule()
    budgeted = has_budget(gl) if budgeted is None else bool(budgeted)
    report = build_report(gl, materiality, budgeted=budgeted)
    comments = line_comments(report, gl, materiality)
    ni = report[report.line == "Net income"].iloc[0]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    have_ytd = ytd is not None and not ytd.empty and budgeted
    elapsed = months or max(1, int(period[5:7]))
    ytd_report = build_report(ytd, materiality, budgeted=True) if have_ytd else None
    # The commentary covers both timeframes, exactly as it does in the workbook.
    comments = line_comments(report, gl, materiality,
                             ytd_report=ytd_report, ytd_gl=ytd)

    ytd_ni = (ytd_report[ytd_report.line == "Net income"].iloc[0]
              if ytd_report is not None else None)
    _cover(prs, entity, period, ni, budgeted, ytd_ni)
    _result(prs, report, gl, period, budgeted, ytd_report)
    _pnl_table(prs, report, comments, period, budgeted, ytd_report)
    if have_ytd:
        _year_to_date(prs, ytd_report, report, elapsed, period, budgeted)
        if fy_budget is not None and not fy_budget.empty:
            _outlook(prs, ytd_report,
                     build_report(fy_budget, materiality, budgeted=True),
                     elapsed, period)
    if budgeted:
        _drivers(prs, gl, materiality, period, budgeted, ytd)
    if analysis_blocks:
        _analysis(prs, analysis_blocks, period)
    if detail is not None and {"department", "entity"} <= set(detail.columns):
        _spend(prs, detail, period, materiality, ytd_detail)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":  # pragma: no cover - manual run
    from .. import ingest
    from ..synthetic_data import (DEFAULT_PERIOD, generate_budget_year,
                                  generate_month, generate_ytd_transactions,
                                  monthly_detail)
    root = Path(__file__).resolve().parents[3]
    period = DEFAULT_PERIOD

    # Roll the demo ledger and plan up to the reporting month.
    cut = ingest.period_key(period)
    txn = generate_ytd_transactions(period)
    bud = generate_budget_year()
    a = (txn[txn.period_no <= cut]
         .groupby(["account_code", "account_name", "category"], as_index=False)
         ["amount_eur"].sum().rename(columns={"amount_eur": "actual"}))
    b = (bud[bud.period_no <= cut]
         .groupby(["account_code", "account_name", "category"], as_index=False)
         [["budget_eur", "prior_eur"]].sum()
         .rename(columns={"budget_eur": "budget", "prior_eur": "prior_year"}))
    ytd_frame = a.merge(b, on=["account_code", "account_name", "category"],
                        how="outer").fillna(0.0)

    fy = (bud.groupby(["account_code", "account_name", "category"], as_index=False)
          [["budget_eur", "prior_eur"]].sum()
          .rename(columns={"budget_eur": "budget", "prior_eur": "prior_year"}))
    fy["actual"] = 0.0

    out = build_pptx_pack(generate_month(period).drop(columns="period"), period,
                          root / "output" / "flux_management_pack.pptx",
                          detail=monthly_detail(period), ytd=ytd_frame,
                          fy_budget=fy,
                          months=txn[txn.period_no <= cut]["period"].nunique())
    print(f"Written: {out}")
